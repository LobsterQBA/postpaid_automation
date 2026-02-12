import os
import io
import copy
import shutil
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from io import BytesIO
from pptx.dml.color import RGBColor
import numpy as np
import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import requests
from pptx.enum.dml import MSO_COLOR_TYPE
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN



def duplicate_slide(prs, slide_index,verbose=False):
    # Create new slide with same layout
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Handle simple shapes with python-pptx API
    for shape in source_slide.shapes:
        shape_type = shape.shape_type
        
        if shape.has_text_frame:
            # Skip shapes with SlideName
            if "SlideName" in shape.text:
                continue
        
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_stream = io.BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(
                image_stream,
                shape.left, shape.top,
                shape.width, shape.height
            )

    # Now use the XML approach to copy all remaining elements
    temp_dir = "temp_pptx"
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)

    # Save the current state to a temporary file
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)

    # Extract the PPTX
    with zipfile.ZipFile(pptx_stream, 'r') as zip_ref:
        zip_ref.extractall(temp_dir)

    # Get paths to important files
    slides_path = os.path.join(temp_dir, "ppt", "slides")
    rels_path = os.path.join(slides_path, "_rels")

    # Get slide files (sort numerically, not alphabetically)
    import re
    def slide_num(filename):
        match = re.search(r'slide(\d+)', filename)
        return int(match.group(1)) if match else 0
    slide_files = sorted([f for f in os.listdir(slides_path) if f.startswith("slide") and f.endswith(".xml")], key=slide_num)
    source_slide_name = slide_files[slide_index]
    new_slide_name = slide_files[-1]  # The new slide we just added will be the last one

    # Get XML content of source slide
    source_slide_xml_path = os.path.join(slides_path, source_slide_name)
    new_slide_xml_path = os.path.join(slides_path, new_slide_name)
    source_slide_rels_path = os.path.join(rels_path, f"{source_slide_name}.rels")
    new_slide_rels_path = os.path.join(rels_path, f"{new_slide_name}.rels")

    # Parse XML of both slides
    source_tree = etree.parse(source_slide_xml_path)
    source_root = source_tree.getroot()
    new_tree = etree.parse(new_slide_xml_path)
    new_root = new_tree.getroot()

    # Parse relationships
    if os.path.exists(source_slide_rels_path) and os.path.exists(new_slide_rels_path):
        source_rels_tree = etree.parse(source_slide_rels_path)
        source_rels_root = source_rels_tree.getroot()
        new_rels_tree = etree.parse(new_slide_rels_path)
        new_rels_root = new_rels_tree.getroot()
        
        # Get existing relationship IDs in new slide
        existing_rIds = {rel.get("Id") for rel in new_rels_root.findall(".//{*}Relationship")}
        
        # Function to generate new relationship IDs
        def gen_new_rId(base="rId"):
            i = 1
            while f"{base}{i}" in existing_rIds:
                i += 1
            new_id = f"{base}{i}"
            existing_rIds.add(new_id)
            return new_id
        
        # Copy relationships for complex elements (tables, charts, etc.)
        old_to_new_rId = {}
        for rel in source_rels_root.findall(".//{*}Relationship"):
            rel_type = rel.get("Type")
            # Skip image relationships as we've already handled those with python-pptx
            if "image" not in rel_type:
                old_rId = rel.get("Id")
                new_rId = gen_new_rId()
                old_to_new_rId[old_rId] = new_rId
                
                # Copy the relationship
                new_rel = etree.SubElement(new_rels_root, "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship")
                new_rel.set("Id", new_rId)
                new_rel.set("Type", rel.get("Type"))
                new_rel.set("Target", rel.get("Target"))
                if rel.get("TargetMode"):
                    new_rel.set("TargetMode", rel.get("TargetMode"))
        
        # Write updated relationships
        new_rels_tree.write(new_slide_rels_path, xml_declaration=True, encoding='UTF-8', standalone="yes")

    # Find all complex shapes in source slide (tables, charts, etc.)
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sp_tree = new_root.find('.//p:spTree', namespaces=ns)

    # Copy complex shapes from source to new slide
    for shape_elem in source_root.findall('.//p:sp', namespaces=ns) + source_root.findall('.//p:graphicFrame', namespaces=ns) + source_root.findall('.//p:grpSp', namespaces=ns):
        # Skip shapes we've already handled with python-pptx (pictures)
        if shape_elem.find('.//p:pic', namespaces=ns) is not None:
            continue
            
        # Skip shapes with "SlideName" text
        txBody = shape_elem.find('.//p:txBody', namespaces=ns)
        if txBody is not None:
            text_content = ''.join(txBody.xpath('.//text()'))
            if "SlideName" in text_content:
                continue
        
        # Copy the shape element
        shape_copy = copy.deepcopy(shape_elem)
        
        # Update relationship IDs in the copied shape
        r_id_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        for elem in shape_copy.iter():
            for attr_name, attr_value in elem.attrib.items():
                if attr_name == r_id_attr and attr_value in old_to_new_rId:
                    elem.attrib[attr_name] = old_to_new_rId[attr_value]
        
        # Add the shape to the new slide
        sp_tree.append(shape_copy)

    # Write the updated slide XML
    new_tree.write(new_slide_xml_path, xml_declaration=True, encoding='UTF-8', standalone="yes")

    # Repackage the presentation
    output_stream = BytesIO()
    with zipfile.ZipFile(output_stream, 'w', zipfile.ZIP_DEFLATED) as zip_out:
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                full_path = os.path.join(root, file)
                arcname = os.path.relpath(full_path, temp_dir)
                zip_out.write(full_path, arcname)

    output_stream.seek(0)
    shutil.rmtree(temp_dir)

    prs = Presentation(output_stream)
    new_slide = prs.slides[-1] 

    find_replace_text(new_slide, "title name: ", "")

    if verbose:
        print_text(new_slide)
        print_tables(new_slide)
        print_images(new_slide)

    
    # Load the modified presentation
    return prs



from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def find_replace_text(slidename, find, replace, verbose=False):
    for shape in slidename.shapes:
        if not shape.has_text_frame:
            continue

        lower_text = shape.text.lower()
        if find.lower() not in lower_text:
            continue

        if verbose:
            print("Original Text:", lower_text)

        text_frame = shape.text_frame
        first_paragraph = text_frame.paragraphs[0]

        if not first_paragraph.runs:
            continue

        old_run = first_paragraph.runs[0]
        font = old_run.font

        # --- Capture old formatting ---
        color = None
        brightness = None
        if font.color is not None:
            if getattr(font.color, "rgb", None):
                color = font.color.rgb
            elif getattr(font.color, "_color", None) and getattr(font.color._color, "theme_color", None):
                color = font.color._color.theme_color
                brightness = font.color.brightness

        # Capture alignment from first paragraph
        old_alignment = first_paragraph.alignment or PP_ALIGN.LEFT

        old_formatting = {
            "name": font.name,
            "bold": font.bold,
            "italic": font.italic,
            "size": font.size,
            "color": color,
            "brightness": brightness,
            "alignment": old_alignment
        }

        # --- Replace text ---
        text_frame.clear()
        text_frame.text = replace

        if verbose:
            print("New Text:", replace)
            print()

        # --- Reapply formatting & alignment ---
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = old_formatting["alignment"]

            for run in paragraph.runs:
                font = run.font
                font.name = old_formatting["name"]
                font.bold = old_formatting["bold"]
                font.italic = old_formatting["italic"]
                font.size = old_formatting["size"]

                if old_formatting["color"] is not None:
                    try:
                        font.color.rgb = old_formatting["color"]
                    except Exception:
                        try:
                            font.color.theme_color = old_formatting["color"]
                            if old_formatting["brightness"] is not None:
                                font.color.brightness = old_formatting["brightness"]
                        except Exception:
                            pass



def print_text(slide):
    print('--- Text in slide ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text
            if "PLACE_TEXT" in text:
                print(text)

def print_tables(slide):
    print('--- Tables in slide ---')
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            first_cell_text = shape.table.cell(0, 0).text.strip()
            if first_cell_text != "":
                print("table name:", first_cell_text)

def print_images(slide):
    print('--- Images in slide ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text
            if "image placeholder" in text:
                print(text)


def create_navigation_screen(prs):
    print('--- Template contains... ---')
    navigation_screen = dict()

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                if "title name" in text:
                    text = shape.text
                    navigation_screen[text.replace('title name: ',"")]=i
                    print(text.replace('title name: ',""), f"- in position of {i}")
    return navigation_screen

def add_data_table_new(slide, tablename, data, header=False):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            first_cell_text = shape.table.cell(0, 0).text.strip()
            if first_cell_text == tablename:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                if not header:
                    rows_new = rows
                else:
                    rows_new = rows - 1

                if data.shape[1] != cols:
                    print('ATTENTION')
                    print(f"Data does not fit in the table, you have {data.shape[1]} columns but the table has {cols}")
                elif data.shape[0] > rows - 1:
                    print('ATTENTION')
                    print(f"Data does not fit in the table, you have {data.shape[0]} rows but the table has {rows-1}")
                else:
                    print(f'Data fits table {tablename}, updating cells with formatting...')
                    for r in range(data.shape[0]):
                        for c in range(data.shape[1]):
                            if not header:
                                cell = table.cell(r + 1, c)  # +1 to skip header
                            else:
                                cell = table.cell(r, c)
                                data_col = pd.DataFrame([data.columns], columns=data.columns)
                                data = pd.concat([data_col, data], ignore_index=True)

                            
                            text_frame = cell.text_frame

                            first_paragraph = text_frame.paragraphs[0]
                            if first_paragraph.runs:
                                old_run = first_paragraph.runs[0]
                                
                                # Store old formatting
                                font = old_run.font

                                try:
                                    color = RGBColor(*font.color.rgb)
                                    brightness= None
                                except:
                                    theme_color = font.color._color.theme_color
                                    color = theme_color
                                    brightness= font.color.brightness   

                                old_formatting = {
                                    "name": font.name,
                                    "bold": font.bold,
                                    "italic": font.italic,
                                    "size": font.size,
                                    "color": color
                                    ,"brightness": brightness
                                }
                            
                            # Clear existing paragraphs
                            p = text_frame.paragraphs[0]
                            p.clear()
                            
                            # Add new run with data
                            run = p.add_run()
                            text_cleaned = str(data.iat[r, c]) if hasattr(data, 'iat') else str(data[r][c])
                            if "[color:" in text_cleaned:
                                if "green" in text_cleaned:
                                    
                                    old_formatting['color'] = RGBColor(0, 176, 80)
                                else:
                                    old_formatting['color'] = RGBColor(255, 0, 0)
                                text_cleaned = text_cleaned.split("[color:")[0]
                            run.text = text_cleaned

                            # Apply original formatting if available
                            font = run.font

                            font.name = old_formatting["name"]
                            font.bold = old_formatting["bold"]
                            font.italic = old_formatting["italic"]
                            font.size = old_formatting["size"]
                            try:
                                font.color.rgb = old_formatting["color"]
                            except:
                                font.color.theme_color = old_formatting["color"]
                                font.color.brightness = old_formatting['brightness']

def format_values_with_color(df, value_column, low_threshold, high_threshold, add_text="", 
                             low_color='red', high_color='green'):
    # Create color-formatted values based on thresholds
    def format_value(val):
        if float(val) >= high_threshold:
            # High value formatting
            color = high_color
            formatted = str(val)+add_text+"[color:{}]".format(color)

        elif float(val) >= low_threshold:
            # Mid value formatting
            formatted = str(val)+add_text
        else:
            # Low value formatting
            color = low_color
            formatted = str(val)+add_text+"[color:{}]".format(color)
        return formatted


    # Apply formatting to each value
    list_val = [format_value(val) for val in df[value_column]]

    return list_val

def add_textbox(slide, text, left, top, width, height, font_size=14, bold=False, color=None, align='left'):
    """Add a text box to a slide"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    
    if color:
        p.font.color.rgb = color
    
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    
    return textbox


def insert_image_to_placeholder(slide, location_name, image_path,verbose=False):
    from pptx.util import Inches
    import os

    # Validate image path
    if not os.path.exists(image_path):
        if verbose:
            print(f"Error: Image file not found at {image_path}")
        return False

    # Find the placeholder shape
    placeholder = None
    search_term = location_name
    for shape in slide.shapes:
        lower_text = shape.text.lower()
        if search_term.lower() in lower_text:
            if verbose:
                print("Original Text:", lower_text)
                placeholder = shape
                break


    # Get placeholder position and size
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    with Image.open(image_path) as img:
        img_width, img_height = img.size
        is_landscape = img_width > img_height

    if is_landscape:
        # Add the image to the slide
        image = slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=width,
                #height=height
            )
    else:
            image = slide.shapes.add_picture(
                image_path,
                left,
                top,
                #width=width,
                height=height
            )

        
        # Optionally remove the placeholder
    sp = placeholder._element
    sp.getparent().remove(sp)
        
    