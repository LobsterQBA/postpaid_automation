"""
PPT Generator - Streamlit Web App
Upload Excel data and generate PowerPoint report
"""

import streamlit as st
import pandas as pd
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from copy import deepcopy

# ========== PAGE CONFIG ==========
st.set_page_config(
    page_title="PPT Report Generator",
    page_icon="📊",
    layout="centered"
)

# ========== CUSTOM CSS (T-Mobile Magenta Theme) ==========
st.markdown("""
<style>
    /* Main theme colors */
    :root {
        --magenta: #E20074;
        --magenta-dark: #B8005C;
        --magenta-light: #FF1493;
    }
    
    /* Header styling */
    .main-header {
        background: linear-gradient(135deg, #E20074 0%, #B8005C 100%);
        padding: 2rem;
        border-radius: 10px;
        text-align: center;
        margin-bottom: 2rem;
    }
    .main-header h1 {
        color: white;
        margin: 0;
        font-size: 2.5rem;
    }
    .main-header p {
        color: rgba(255,255,255,0.9);
        margin: 0.5rem 0 0 0;
    }
    
    /* Button styling */
    .stButton > button {
        background: linear-gradient(135deg, #E20074 0%, #B8005C 100%);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        font-size: 1.1rem;
        font-weight: 600;
        border-radius: 8px;
        width: 100%;
        transition: all 0.3s ease;
    }
    .stButton > button:hover {
        background: linear-gradient(135deg, #FF1493 0%, #E20074 100%);
        box-shadow: 0 4px 15px rgba(226, 0, 116, 0.4);
    }
    
    /* Upload area styling */
    .uploadedFile {
        border: 2px dashed #E20074 !important;
    }
    
    /* Success message */
    .success-box {
        background: linear-gradient(135deg, #00C853 0%, #00E676 100%);
        padding: 1rem;
        border-radius: 8px;
        text-align: center;
        color: white;
        font-weight: 600;
    }
    
    /* Info box */
    .info-box {
        background: #f8f9fa;
        border-left: 4px solid #E20074;
        padding: 1rem;
        border-radius: 0 8px 8px 0;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

# ========== HEADER ==========
st.markdown("""
<div class="main-header">
    <h1>📊 PPT Report Generator</h1>
    <p>Upload your Excel data and generate a PowerPoint report instantly</p>
</div>
""", unsafe_allow_html=True)

# ========== HELPER FUNCTIONS ==========
def format_number(x):
    """Format numbers with comma separators"""
    if pd.isna(x) or x == "":
        return ""
    try:
        x = int(float(x))
        return f"{x:,}"
    except:
        return str(x)

def format_percent(x):
    """Format as percentage"""
    if pd.isna(x) or x == "":
        return ""
    try:
        x = float(x)
        if x < 1:
            return f"{x*100:.2f}%"
        else:
            return f"{x:.2f}%"
    except:
        return str(x)

def process_em_data(em_data):
    """Process EM data for PPT table"""
    df = em_data.copy()
    df['Open Rate'] = df['Unique Opens'] / df['Deliveries']
    df['Click Rate'] = df['Unique Clicks'] / df['Deliveries']
    
    table_df = df[['Group', 'Offer type', 'Deliveries', 'Unique Opens', 'Unique Clicks', 'Open Rate', 'Click Rate', 'Subject Line']].copy()
    table_df.columns = ['Group', 'Offer Type', 'Deliveries', 'Opens', 'Clicks', 'Open Rate', 'Click Rate', 'Creative Message']
    
    table_df['Deliveries'] = table_df['Deliveries'].apply(format_number)
    table_df['Opens'] = table_df['Opens'].apply(format_number)
    table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    table_df['Open Rate'] = table_df['Open Rate'].apply(format_percent)
    table_df['Click Rate'] = table_df['Click Rate'].apply(format_percent)
    
    return table_df

def process_sms_data(sms_data):
    """Process SMS data for PPT table"""
    df = sms_data.copy()
    table_df = df[['Touch', 'Channel', 'Group', 'Sent (Est.)', 'Clicks', 'CTR', 'Creative Message']].copy()
    table_df.columns = ['Touch', 'Channel', 'Group', 'Sent', 'Clicks', 'CTR', 'Creative Message']
    
    table_df['Sent'] = table_df['Sent'].apply(format_number)
    table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    table_df['CTR'] = table_df['CTR'].apply(format_percent)
    
    return table_df

def process_sl_testing_data(sl_data):
    """Process SL Testing data"""
    df = sl_data[sl_data['Row Labels'] != 'Grand Total'].copy()
    
    df = df.rename(columns={
        'Row Labels': 'Variant',
        'Sum of Unique People Clicks [v8]': 'Clicks',
        'Sum of People Opens [v8]': 'Opens',
        'Sum of Unique People Delivered [v8]': 'Delivered',
        'Open Rate': 'Open Rate',
        'Click Rate': 'Click Rate',
        'Subject Line': 'Subject Line'
    })
    
    table_df = df[['Variant', 'Delivered', 'Opens', 'Open Rate', 'Clicks', 'Click Rate', 'Subject Line']].copy()
    
    table_df['Delivered'] = table_df['Delivered'].apply(format_number)
    table_df['Opens'] = table_df['Opens'].apply(format_number)
    table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    table_df['Open Rate'] = table_df['Open Rate'].apply(format_percent)
    table_df['Click Rate'] = table_df['Click Rate'].apply(format_percent)
    
    return table_df


def add_data_to_table(slide, data, font_size=9):
    """Add data to table, dynamically adding rows if needed"""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            cols = len(table.columns)
            needed_rows = data.shape[0] + 1
            
            while len(table.rows) < needed_rows:
                tbl = table._tbl
                tr = tbl.tr_lst[-1]
                new_tr = deepcopy(tr)
                tbl.append(new_tr)
            
            for c, col_name in enumerate(data.columns):
                if c < cols:
                    cell = table.cell(0, c)
                    p = cell.text_frame.paragraphs[0]
                    p.clear()
                    run = p.add_run()
                    run.text = col_name
                    run.font.size = Pt(font_size + 1)
                    run.font.bold = True
            
            for r in range(data.shape[0]):
                for c in range(min(data.shape[1], cols)):
                    cell = table.cell(r + 1, c)
                    p = cell.text_frame.paragraphs[0]
                    text = str(data.iloc[r, c]) if pd.notna(data.iloc[r, c]) else ""
                    p.clear()
                    run = p.add_run()
                    run.text = text
                    run.font.size = Pt(font_size)
            
            for c in range(data.shape[1], cols):
                for r in range(len(table.rows)):
                    cell = table.cell(r, c)
                    cell.text_frame.paragraphs[0].clear()
            return

def find_replace_text(slide, find, replace):
    """Find and replace text in slide"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if find.lower() in shape.text.lower():
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if find.lower() in run.text.lower():
                        run.text = run.text.replace(find, replace).replace(find.upper(), replace)

def duplicate_slide_simple(prs, slide_index):
    """Simple slide duplication"""
    source = prs.slides[slide_index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)
    
    for shape in source.shapes:
        if shape.has_text_frame:
            if "SlideName" in shape.text or "title name" in shape.text:
                continue
    
    return new_slide

def create_navigation_screen(prs):
    """Create navigation dictionary from template"""
    navigation = {}
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "title name" in shape.text:
                    name = shape.text.replace('title name: ', '')
                    navigation[name] = i
    return navigation

def generate_ppt(em_data, sms_data, sl_data, template_bytes):
    """Generate PPT from data"""
    import functions as DS
    
    prs = Presentation(io.BytesIO(template_bytes))
    nav = DS.create_navigation_screen(prs)
    
    # Title slide
    prs = DS.duplicate_slide(prs, nav['titleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Weekly Report', verbose=False)
    
    # Email Performance
    prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Performance', verbose=False)
    
    em_table = process_em_data(em_data)
    prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Performance Overview', verbose=False)
    add_data_to_table(prs.slides[-1], em_table)
    
    # SMS/RBM Performance
    prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS/RBM Performance', verbose=False)
    
    sms_table = process_sms_data(sms_data)
    prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS/RBM Performance Overview', verbose=False)
    add_data_to_table(prs.slides[-1], sms_table)
    
    # Subject Line Testing
    prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Subject Line Testing', verbose=False)
    
    sl_table = process_sl_testing_data(sl_data)
    prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Subject Line Testing Results', verbose=False)
    add_data_to_table(prs.slides[-1], sl_table)
    
    # Cleanup template slides
    for i in sorted(nav.values(), reverse=True):
        rId = prs.slides._sldIdLst[i]
        prs.slides._sldIdLst.remove(rId)
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ========== MAIN UI ==========
st.markdown('<div class="info-box">📁 <strong>Required format:</strong> Excel file with 3 tabs: <code>EM</code>, <code>SMS</code>, <code>SL Testing</code></div>', unsafe_allow_html=True)

# File upload
uploaded_file = st.file_uploader(
    "Drop your Excel data file here",
    type=['xlsx', 'xls'],
    help="Upload an Excel file with EM, SMS, and SL Testing tabs"
)

if uploaded_file is not None:
    try:
        # Read data
        em_data = pd.read_excel(uploaded_file, sheet_name='EM')
        uploaded_file.seek(0)
        sms_data = pd.read_excel(uploaded_file, sheet_name='SMS')
        uploaded_file.seek(0)
        sl_data = pd.read_excel(uploaded_file, sheet_name='SL Testing')
        
        # Show preview
        st.success(f"✅ File loaded successfully!")
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("EM Rows", len(em_data))
        with col2:
            st.metric("SMS Rows", len(sms_data))
        with col3:
            st.metric("SL Testing Rows", len(sl_data) - 1)  # -1 for Grand Total
        
        # Preview expanders
        with st.expander("📧 Preview EM Data"):
            st.dataframe(em_data.head(), use_container_width=True)
        
        with st.expander("💬 Preview SMS Data"):
            st.dataframe(sms_data.head(), use_container_width=True)
        
        with st.expander("🔬 Preview SL Testing Data"):
            st.dataframe(sl_data.head(), use_container_width=True)
        
        st.markdown("---")
        
        # Generate button
        if st.button("🚀 Generate PPT Report", use_container_width=True):
            with st.spinner("Generating your PPT report..."):
                # Load template
                with open('template.pptx', 'rb') as f:
                    template_bytes = f.read()
                
                # Generate PPT
                ppt_output = generate_ppt(em_data, sms_data, sl_data, template_bytes)
                
                st.markdown('<div class="success-box">✅ PPT Report Generated Successfully!</div>', unsafe_allow_html=True)
                
                # Download button
                st.download_button(
                    label="📥 Download PPT Report",
                    data=ppt_output,
                    file_name="weekly_report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
    
    except Exception as e:
        st.error(f"❌ Error: {str(e)}")
        st.info("Please make sure your Excel file has the correct tabs: EM, SMS, SL Testing")

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #888; font-size: 0.9rem;">
    Made with ❤️ using Streamlit | T-Mobile Marketing Analytics
</div>
""", unsafe_allow_html=True)
