"""
Campaign PPT Generator - Streamlit Web App
Uses unified "Clean Data Sheet" format matching Data Repository standards
"""

import streamlit as st
import pandas as pd
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
import os
from pulse_check_docx import build_pulse_check_docx

# ========== PAGE CONFIG ==========
st.set_page_config(
    page_title="Campaign PPT Generator",
    page_icon="📊",
    layout="wide"
)

# ========== CUSTOM CSS (T-Mobile Magenta Theme) ==========
st.markdown("""
<style>
    :root {
        --magenta: #E20074;
        --magenta-dark: #B8005C;
    }
    .main-header {
        background: linear-gradient(135deg, #E20074 0%, #B8005C 100%);
        padding: 2rem;
        border-radius: 10px;
        text-align: center;
        margin-bottom: 2rem;
    }
    .main-header h1 { color: white; margin: 0; font-size: 2.5rem; }
    .main-header p { color: rgba(255,255,255,0.9); margin: 0.5rem 0 0 0; }
    .stButton > button {
        background: linear-gradient(135deg, #E20074 0%, #B8005C 100%);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        font-size: 1.1rem;
        font-weight: 600;
        border-radius: 8px;
        width: 100%;
    }
    .stButton > button:hover {
        background: linear-gradient(135deg, #FF1493 0%, #E20074 100%);
        box-shadow: 0 4px 15px rgba(226, 0, 116, 0.4);
    }
    .info-box {
        background: #f8f9fa;
        border-left: 4px solid #E20074;
        padding: 1rem;
        border-radius: 0 8px 8px 0;
        margin: 1rem 0;
    }
    .success-box {
        background: linear-gradient(135deg, #00C853 0%, #00E676 100%);
        padding: 1rem;
        border-radius: 8px;
        text-align: center;
        color: white;
        font-weight: 600;
    }
</style>
""", unsafe_allow_html=True)

# ========== HEADER ==========
st.markdown("""
<div class="main-header">
    <h1>📊 Campaign PPT Generator</h1>
    <p>Generate standardized Campaign Performance Reports using Data Repository format</p>
</div>
""", unsafe_allow_html=True)


# ========== HELPER FUNCTIONS ==========
def format_number(x):
    """Format numbers with comma separators"""
    if pd.isna(x) or x == "":
        return ""
    try:
        x = int(float(x))
        return f"{x:,}"
    except:
        return str(x)


def format_percent(x):
    """Format as percentage"""
    if pd.isna(x) or x == "":
        return ""
    try:
        x = float(x)
        if x < 1:
            return f"{x*100:.2f}%"
        else:
            return f"{x:.2f}%"
    except:
        return str(x)


def truncate_text(text, max_length=50):
    """Truncate long text with ellipsis"""
    if pd.isna(text):
        return ""
    text = str(text)
    if len(text) > max_length:
        return text[:max_length-3] + "..."
    return text


def is_column_empty(df, col_name):
    """Check if a column is completely empty or all NaN"""
    if col_name not in df.columns:
        return True
    return df[col_name].isna().all() or (df[col_name].astype(str).str.strip() == '').all()


def get_audience_columns(df):
    """Get list of non-empty Audience Details columns"""
    audience_cols = []
    for col in ['Audience Details 1', 'Audience Details 2', 'Audience Details 3']:
        if col in df.columns and not is_column_empty(df, col):
            audience_cols.append(col)
    return audience_cols


def auto_fit_table_columns(table, min_width=0.5, max_width=2.5):
    """Adjust column widths based on content length"""
    from pptx.util import Inches
    num_cols = len(table.columns)
    num_rows = len(table.rows)
    
    for col_idx in range(num_cols):
        max_len = 0
        for row_idx in range(num_rows):
            cell = table.cell(row_idx, col_idx)
            text_len = len(str(cell.text)) if cell.text else 0
            max_len = max(max_len, text_len)
        # Calculate width proportional to content, with min/max bounds
        width = min(max(max_len * 0.1, min_width), max_width)
        table.columns[col_idx].width = Inches(width)


def add_data_to_table(slide, data, font_size=9):
    """Add data to table, dynamically adding rows if needed"""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            cols = len(table.columns)
            needed_rows = data.shape[0] + 1  # +1 for header
            
            # Add rows if needed
            while len(table.rows) < needed_rows:
                tbl = table._tbl
                tr = tbl.tr_lst[-1]
                new_tr = deepcopy(tr)
                tbl.append(new_tr)
            
            # Update headers (row 0)
            for c, col_name in enumerate(data.columns):
                if c < cols:
                    cell = table.cell(0, c)
                    p = cell.text_frame.paragraphs[0]
                    p.clear()
                    run = p.add_run()
                    run.text = col_name
                    run.font.size = Pt(font_size + 1)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White text
            
            # Apply Magenta background to header row
            for i in range(cols):
                cell = table.cell(0, i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(226, 0, 116)  # Magenta #E20074
            
            # Update data rows (row 1+)
            for r in range(data.shape[0]):
                for c in range(min(data.shape[1], cols)):
                    cell = table.cell(r + 1, c)
                    p = cell.text_frame.paragraphs[0]
                    text = str(data.iloc[r, c]) if pd.notna(data.iloc[r, c]) else ""
                    p.clear()
                    run = p.add_run()
                    run.text = text
                    run.font.size = Pt(font_size)
            
            # Clear unused columns
            for c in range(data.shape[1], cols):
                for r in range(len(table.rows)):
                    cell = table.cell(r, c)
                    cell.text_frame.paragraphs[0].clear()
            
            # Auto-fit column widths based on content
            auto_fit_table_columns(table)
            return


def find_replace_text(slide, find, replace):
    """Find and replace text in slide, preserving formatting"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if find.lower() in shape.text.lower():
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if find.lower() in run.text.lower():
                        run.text = run.text.replace(find, replace).replace(find.upper(), replace)


# ========== DATA PROCESSING FUNCTIONS ==========

def process_email_summary(em_data):
    """
    Create Email Summary table aggregated by Touch
    Columns: Touch | Deliveries | Open Rate | CTR
    """
    if em_data is None or len(em_data) == 0:
        return pd.DataFrame(columns=['Touch', 'Deliveries', 'Open Rate', 'CTR'])
    
    # Group by Touch and aggregate
    summary = em_data.groupby('Touch').agg({
        'Deliveries': 'sum',
        'Unique Opens': 'sum',
        'Unique Clicks': 'sum'
    }).reset_index()
    
    # Calculate rates
    summary['Open Rate'] = summary['Unique Opens'] / summary['Deliveries']
    summary['CTR'] = summary['Unique Clicks'] / summary['Deliveries']
    
    # Format for display
    result = pd.DataFrame({
        'Touch': summary['Touch'],
        'Deliveries': summary['Deliveries'].apply(format_number),
        'Open Rate': summary['Open Rate'].apply(format_percent),
        'CTR': summary['CTR'].apply(format_percent)
    })
    
    return result


def process_email_summary_total(em_data):
    """
    Create Email Summary table aggregated across all Touches
    Columns: Segment | Deliveries | Open Rate | CTR
    """
    if em_data is None or len(em_data) == 0:
        return pd.DataFrame(columns=['Segment', 'Deliveries', 'Open Rate', 'CTR'])

    total_deliveries = em_data['Deliveries'].sum()
    total_opens = em_data['Unique Opens'].sum()
    total_clicks = em_data['Unique Clicks'].sum()

    open_rate = total_opens / total_deliveries if total_deliveries else 0
    ctr = total_clicks / total_deliveries if total_deliveries else 0

    result = pd.DataFrame([{
        'Segment': 'Total',
        'Deliveries': format_number(total_deliveries),
        'Open Rate': format_percent(open_rate),
        'CTR': format_percent(ctr)
    }])

    return result


def process_sms_summary(sms_data):
    """
    Create SMS Summary table aggregated by Touch
    Columns: Touch | Deliveries | CTR
    """
    if sms_data is None or len(sms_data) == 0:
        return pd.DataFrame(columns=['Touch', 'Deliveries', 'CTR'])
    
    # Group by Touch and aggregate
    summary = sms_data.groupby('Touch').agg({
        'Deliveries': 'sum',
        'PBI Unique Clicks': 'sum'
    }).reset_index()
    
    # Calculate CTR
    summary['CTR'] = summary['PBI Unique Clicks'] / summary['Deliveries']
    
    # Format for display
    result = pd.DataFrame({
        'Touch': summary['Touch'],
        'Deliveries': summary['Deliveries'].apply(format_number),
        'CTR': summary['CTR'].apply(format_percent)
    })
    
    return result


def process_sms_summary_total(sms_data):
    """
    Create SMS Summary table aggregated across all Touches
    Columns: Segment | Deliveries | CTR
    """
    if sms_data is None or len(sms_data) == 0:
        return pd.DataFrame(columns=['Segment', 'Deliveries', 'CTR'])

    total_deliveries = sms_data['Deliveries'].sum()
    total_clicks = sms_data['PBI Unique Clicks'].sum() if 'PBI Unique Clicks' in sms_data.columns else 0
    ctr = total_clicks / total_deliveries if total_deliveries else 0

    result = pd.DataFrame([{
        'Segment': 'Total',
        'Deliveries': format_number(total_deliveries),
        'CTR': format_percent(ctr)
    }])

    return result

def process_em_data(em_data):
    """
    Process EM data for Email Performance Overview slide
    Table columns (in order):
    1. Touch
    2. OS
    3. Cohort
    4. Audience Details 1 (if has data)
    5. Deliveries
    6. Unique Opens
    7. Unique Clicks
    8. Open Rate
    9. Click Rate
    
    NO Subject Line column!
    """
    df = em_data.copy()
    
    # Calculate rates
    df['Open Rate'] = df['Unique Opens'] / df['Deliveries']
    df['Click Rate'] = df['Unique Clicks'] / df['Deliveries']
    
    # Build columns list in exact order specified
    display_cols = []
    
    # Add Touch, OS, Cohort first
    for col in ['Touch', 'OS', 'Cohort']:
        if col in df.columns:
            display_cols.append(col)
    
    # Add Audience Details 1 only (if has data)
    if 'Audience Details 1' in df.columns and not is_column_empty(df, 'Audience Details 1'):
        display_cols.append('Audience Details 1')
    
    # Add metrics in exact order: Deliveries, Unique Opens, Unique Clicks, Open Rate, Click Rate
    display_cols.extend(['Deliveries', 'Unique Opens', 'Unique Clicks', 'Open Rate', 'Click Rate'])
    
    # Filter to existing columns
    display_cols = [c for c in display_cols if c in df.columns]
    
    table_df = df[display_cols].copy()
    
    # Format numbers
    if 'Deliveries' in table_df.columns:
        table_df['Deliveries'] = table_df['Deliveries'].apply(format_number)
    if 'Unique Opens' in table_df.columns:
        table_df['Unique Opens'] = table_df['Unique Opens'].apply(format_number)
    if 'Unique Clicks' in table_df.columns:
        table_df['Unique Clicks'] = table_df['Unique Clicks'].apply(format_number)
    table_df['Open Rate'] = table_df['Open Rate'].apply(format_percent)
    table_df['Click Rate'] = table_df['Click Rate'].apply(format_percent)
    
    return table_df


def process_sms_data(sms_data, has_testing=False):
    """
    Process SMS data for SMS Performance Overview slide
    
    Columns (in order):
    1. Touch
    2. OS
    3. Cohort
    4. Audience Details 1/2/3 (if has data)
    5. Creative
    6. Deliveries
    7. Clicks (PBI Unique Clicks)
    8. CTR
    
    Note: SMS Testing Variant is handled separately in process_sms_testing_data()
    """
    df = sms_data.copy()
    
    # Calculate Delivery Rate (keep for potential use)
    df['Delivery Rate'] = df['Deliveries'] / df['Sends']
    
    # Rename columns to display names
    if 'PBI CTR' in df.columns:
        df['CTR'] = df['PBI CTR']
    elif 'PBI  CTR' in df.columns:  # Handle double-space variant
        df['CTR'] = df['PBI  CTR']
    
    if 'PBI Unique Clicks' in df.columns:
        df['Clicks'] = df['PBI Unique Clicks']
    
    # Build columns list
    display_cols = []
    
    # 1. Touch, OS, Cohort first
    for col in ['Touch', 'OS', 'Cohort']:
        if col in df.columns:
            display_cols.append(col)
    
    # 2. Audience Details columns (only if they have data)
    for col in ['Audience Details 1', 'Audience Details 2', 'Audience Details 3']:
        if col in df.columns and not is_column_empty(df, col):
            display_cols.append(col)
    
    # 3. Creative / Message Body
    if 'Creative' in df.columns:
        display_cols.append('Creative')
    
    # 4. Deliveries
    display_cols.append('Deliveries')
    
    # 5. Clicks
    display_cols.append('Clicks')
    
    # 6. CTR
    display_cols.append('CTR')
    
    # Filter to existing columns
    display_cols = [c for c in display_cols if c in df.columns]
    
    table_df = df[display_cols].copy()
    
    # Truncate Creative column to fit
    if 'Creative' in table_df.columns:
        table_df['Creative'] = table_df['Creative'].apply(lambda x: truncate_text(x, 50))
    
    # Format numbers
    if 'Deliveries' in table_df.columns:
        table_df['Deliveries'] = table_df['Deliveries'].apply(format_number)
    if 'Clicks' in table_df.columns:
        table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    if 'CTR' in table_df.columns:
        table_df['CTR'] = table_df['CTR'].apply(format_percent)
    
    return table_df


def process_sl_data(sl_data, em_data=None):
    """
    Process Subject Line Testing data
    Table columns (in order):
    1. Subject Line
    2. Cohort
    3. OS (Operating System)
    4. Deliveries
    5. Open Rate (OR)
    6. Click Rate (CTR)
    """
    df = sl_data.copy()
    
    # Build the table with exact columns
    display_cols = ['Subject Line']
    
    # Add Cohort column
    if 'Cohort' in df.columns:
        display_cols.append('Cohort')
    
    # Add OS column (from Operating System)
    if 'Operating System' in df.columns:
        df['OS'] = df['Operating System']
        display_cols.append('OS')

    # Deliveries (support common variants)
    if 'Delivered' in df.columns:
        df['Deliveries'] = df['Delivered']
        display_cols.append('Deliveries')
    elif 'Deliveries' in df.columns:
        display_cols.append('Deliveries')
    elif 'Sum of Unique People Delivered [v8]' in df.columns:
        df['Deliveries'] = df['Sum of Unique People Delivered [v8]']
        display_cols.append('Deliveries')
        
    # Add Winner/Loser column (empty for manual entry)
    df['Winner/Loser'] = ''
    display_cols.append('Winner/Loser')
    
    # Open Rate and Click Rate
    if 'OR' in df.columns:
        df['Open Rate'] = df['OR']
        display_cols.append('Open Rate')
    
    if 'CTR' in df.columns:
        df['Click Rate'] = df['CTR']
        display_cols.append('Click Rate')
    
    # Filter to existing columns (Winner/Loser is newly created so it will exist)
    display_cols = [c for c in display_cols if c in df.columns]
    
    table_df = df[display_cols].copy()
    
    # Truncate subject lines
    if 'Subject Line' in table_df.columns:
        table_df['Subject Line'] = table_df['Subject Line'].apply(lambda x: truncate_text(x, 45))
    
    # Format percentages
    if 'Deliveries' in table_df.columns:
        table_df['Deliveries'] = table_df['Deliveries'].apply(format_number)
    if 'Open Rate' in table_df.columns:
        table_df['Open Rate'] = table_df['Open Rate'].apply(format_percent)
    if 'Click Rate' in table_df.columns:
        table_df['Click Rate'] = table_df['Click Rate'].apply(format_percent)
    
    return table_df


def process_sms_testing_data(sms_data):
    """
    Process SMS Testing data - only called when SMS Testing Variant has data
    Table columns:
    1. SMS Testing Variant
    2. Touch
    3. Creative
    4. Deliveries
    5. Clicks
    6. CTR
    """
    df = sms_data.copy()
    
    # Rename columns to display names
    if 'PBI CTR' in df.columns:
        df['CTR'] = df['PBI CTR']
    elif 'PBI  CTR' in df.columns:
        df['CTR'] = df['PBI  CTR']
    
    if 'PBI Unique Clicks' in df.columns:
        df['Clicks'] = df['PBI Unique Clicks']
    
    # Build columns list - SMS Testing Variant first
    display_cols = ['SMS Testing Variant']
    
    # Add Touch
    if 'Touch' in df.columns:
        display_cols.append('Touch')
    
    # Add Creative
    if 'Creative' in df.columns:
        display_cols.append('Creative')
    
    # Add metrics
    display_cols.extend(['Deliveries', 'Clicks', 'CTR'])
    
    # Filter to existing columns
    display_cols = [c for c in display_cols if c in df.columns]
    
    table_df = df[display_cols].copy()
    
    # Truncate Creative column
    if 'Creative' in table_df.columns:
        table_df['Creative'] = table_df['Creative'].apply(lambda x: truncate_text(x, 40))
    
    # Format numbers
    if 'Deliveries' in table_df.columns:
        table_df['Deliveries'] = table_df['Deliveries'].apply(format_number)
    if 'Clicks' in table_df.columns:
        table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    if 'CTR' in table_df.columns:
        table_df['CTR'] = table_df['CTR'].apply(format_percent)
    
    return table_df


def check_sms_has_testing(sms_data):
    """Check if SMS data has testing variants"""
    if 'SMS Testing Variant' not in sms_data.columns:
        return False
    return not is_column_empty(sms_data, 'SMS Testing Variant')


# ========== PPT GENERATION ==========
def generate_campaign_ppt(em_data, sms_data, sl_data, campaign_name, template_bytes):
    """Generate PPT from campaign data"""
    import functions as DS
    
    prs = Presentation(io.BytesIO(template_bytes))
    nav = DS.create_navigation_screen(prs)
    
    # === Title Slide ===
    prs = DS.duplicate_slide(prs, nav['titleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', campaign_name, verbose=False)
    
    # === Campaign Summary Section ===
    prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Campaign Summary', verbose=False)
    
    # Email Summary table (Total)
    if em_data is not None and len(em_data) > 0:
        email_summary_total = process_email_summary_total(em_data)
        prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Summary (Total)', verbose=False)
        add_data_to_table(prs.slides[-1], email_summary_total)

    # Email Summary table (by Touch)
    if em_data is not None and len(em_data) > 0:
        email_summary = process_email_summary(em_data)
        prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Summary by Touch', verbose=False)
        add_data_to_table(prs.slides[-1], email_summary)
    
    # SMS Summary table (Total)
    if sms_data is not None and len(sms_data) > 0:
        sms_summary_total = process_sms_summary_total(sms_data)
        prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS Summary (Total)', verbose=False)
        add_data_to_table(prs.slides[-1], sms_summary_total)

    # SMS Summary table (by Touch)
    if sms_data is not None and len(sms_data) > 0:
        sms_summary = process_sms_summary(sms_data)
        prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS Summary by Touch', verbose=False)
        add_data_to_table(prs.slides[-1], sms_summary)
    
    # === Email High-Level Section ===
    if em_data is not None and len(em_data) > 0:
        prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email High-Level Results', verbose=False)
        
        em_table = process_em_data(em_data)
        prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Performance Overview', verbose=False)
        add_data_to_table(prs.slides[-1], em_table)
    
    # === SMS High-Level Section ===
    if sms_data is not None and len(sms_data) > 0:
        prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS High-Level Results', verbose=False)
        
        # SMS Detailed Data table (full data)
        sms_table = process_sms_data(sms_data)
        prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS Performance Overview', verbose=False)
        add_data_to_table(prs.slides[-1], sms_table)
    
    # === Subject Line Testing Section (only if SLs sheet has data) ===
    if sl_data is not None and len(sl_data) > 0:
        prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Subject Line Testing', verbose=False)
        
        sl_table = process_sl_data(sl_data)
        prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Subject Line Testing Results', verbose=False)
        add_data_to_table(prs.slides[-1], sl_table.head(15))  # Limit rows for readability
    
    # === SMS Testing Section (only if SMS Testing Variant has data) ===
    if sms_data is not None and len(sms_data) > 0:
        has_testing = check_sms_has_testing(sms_data)
        if has_testing:
            prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
            DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS Testing', verbose=False)
            
            sms_testing_table = process_sms_testing_data(sms_data)
            prs = DS.duplicate_slide(prs, nav['smsdataslide'], verbose=False)
            DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS Testing Results', verbose=False)
            add_data_to_table(prs.slides[-1], sms_testing_table)
    
    # === Cleanup template slides ===
    for i in sorted(nav.values(), reverse=True):
        rId = prs.slides._sldIdLst[i]
        prs.slides._sldIdLst.remove(rId)
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _metric_weighted_widths(columns):
    metric_headers = {
        "Deliveries", "Unique Opens", "Unique Clicks", "Open Rate", "Click Rate",
        "Clicks", "CTR", "OR", "Sent", "Sent (Est.)"
    }
    weights = []
    for col in columns:
        if col in {"Creative", "Subject Line"}:
            weights.append(3)
        elif col in {"Touch", "OS", "Cohort", "Audience Details 1", "Audience Details 2", "Audience Details 3", "SMS Testing Variant"}:
            weights.append(2)
        elif col in metric_headers:
            weights.append(1)
        else:
            weights.append(2)
    total = sum(weights) if weights else 1
    return [w / total for w in weights]


def _pulse_check_sections(em_data, sms_data, sl_data):
    sections = []

    if em_data is not None and len(em_data) > 0:
        em_total = process_email_summary_total(em_data)
        sections.append({
            "header": "Email Summary (Total)",
            "data": em_total,
            "widths": [0.30, 0.23, 0.23, 0.24]
        })
        em_touch = process_email_summary(em_data)
        sections.append({
            "header": "Email Summary by Touch",
            "data": em_touch,
            "widths": [0.30, 0.23, 0.23, 0.24]
        })

    if sms_data is not None and len(sms_data) > 0:
        sms_total = process_sms_summary_total(sms_data)
        sections.append({
            "header": "SMS Summary (Total)",
            "data": sms_total,
            "widths": [0.40, 0.30, 0.30]
        })
        sms_touch = process_sms_summary(sms_data)
        sections.append({
            "header": "SMS Summary by Touch",
            "data": sms_touch,
            "widths": [0.40, 0.30, 0.30]
        })

    if em_data is not None and len(em_data) > 0:
        em_perf = process_em_data(em_data)
        sections.append({
            "header": "Email Performance Overview",
            "data": em_perf,
            "widths": _metric_weighted_widths(em_perf.columns)
        })

    if sms_data is not None and len(sms_data) > 0:
        sms_perf = process_sms_data(sms_data)
        sections.append({
            "header": "SMS Performance Overview",
            "data": sms_perf,
            "widths": _metric_weighted_widths(sms_perf.columns)
        })

    if sl_data is not None and len(sl_data) > 0:
        sl_table = process_sl_data(sl_data)
        sections.append({
            "header": "Subject Line Testing Results",
            "data": sl_table,
            "widths": _metric_weighted_widths(sl_table.columns)
        })

    if sms_data is not None and len(sms_data) > 0:
        has_testing = check_sms_has_testing(sms_data)
        if has_testing:
            sms_testing = process_sms_testing_data(sms_data)
            sections.append({
                "header": "SMS Testing Results",
                "data": sms_testing,
                "widths": _metric_weighted_widths(sms_testing.columns)
            })

    return sections


# ========== MAIN UI ==========

# File upload section
st.markdown('<div class="info-box">📁 <strong>Data Format:</strong> Excel file with tabs: <code>EM</code>, <code>SMS</code>, <code>SLs</code> (Data Repository standard)</div>', unsafe_allow_html=True)

# Option to use default Sample Data or upload custom file
use_sample = st.checkbox("Use PPT_Sample_Input_Data.xlsx from project folder", value=True)

if use_sample:
    data_file_path = os.path.join(os.path.dirname(__file__), 'PPT_Sample_Input_Data.xlsx')
    if os.path.exists(data_file_path):
        uploaded_file = data_file_path
        st.success(f"✅ Using: {data_file_path}")
    else:
        st.error("❌ PPT_Sample_Input_Data.xlsx not found in project folder")
        uploaded_file = None
else:
    uploaded_file = st.file_uploader(
        "Drop your Excel data file here",
        type=['xlsx', 'xls'],
        help="Upload an Excel file with EM, SMS, and SLs tabs"
    )

if uploaded_file is not None:
    try:
        # Read data
        if isinstance(uploaded_file, str):
            em_data = pd.read_excel(uploaded_file, sheet_name='EM')
            sms_data = pd.read_excel(uploaded_file, sheet_name='SMS')
            sl_data = pd.read_excel(uploaded_file, sheet_name='SLs')
        else:
            em_data = pd.read_excel(uploaded_file, sheet_name='EM')
            uploaded_file.seek(0)
            sms_data = pd.read_excel(uploaded_file, sheet_name='SMS')
            uploaded_file.seek(0)
            sl_data = pd.read_excel(uploaded_file, sheet_name='SLs')
        
        st.success("✅ Data loaded successfully!")
        
        # Show data summary
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("EM Rows", len(em_data))
        with col2:
            st.metric("SMS Rows", len(sms_data))
        with col3:
            st.metric("SL Testing Rows", len(sl_data))
        
        # Campaign selection
        all_campaigns = set()
        all_campaigns.update(em_data['Campaign'].unique().tolist())
        all_campaigns.update(sms_data['Campaign'].unique().tolist())
        all_campaigns = sorted(list(all_campaigns))
        
        st.markdown("---")
        st.subheader("📋 Campaign Selection")
        
        selected_campaign = st.selectbox(
            "Select Campaign to Generate Report",
            options=["All Campaigns"] + all_campaigns,
            help="Choose a specific campaign or generate for all"
        )
        
        # Filter data by campaign
        if selected_campaign != "All Campaigns":
            em_filtered = em_data[em_data['Campaign'] == selected_campaign]
            sms_filtered = sms_data[sms_data['Campaign'] == selected_campaign]
            sl_filtered = sl_data[sl_data['Campaign'] == selected_campaign]
            campaign_name = selected_campaign
        else:
            em_filtered = em_data
            sms_filtered = sms_data
            sl_filtered = sl_data
            campaign_name = "Campaign Performance Report"
        
        # Show filtered data preview
        with st.expander("📧 Preview Email Data"):
            st.dataframe(em_filtered.head(10), use_container_width=True)
        
        with st.expander("💬 Preview SMS Data"):
            st.dataframe(sms_filtered.head(10), use_container_width=True)
        
        with st.expander("🔬 Preview Subject Line Data"):
            st.dataframe(sl_filtered.head(10), use_container_width=True)
        
        # SMS Testing status
        has_sms_testing = check_sms_has_testing(sms_filtered)
        if has_sms_testing:
            st.info("📊 SMS Testing Variant detected - will generate Testing layout")
        else:
            st.info("📊 No SMS Testing Variant - will generate standard layout")
        
        st.markdown("---")
        st.subheader("📄 Output Style")
        report_style = st.radio(
            "Choose output format",
            options=[
                "Campaign Report (PPT)",
                "Pulse Check (Word .docx)"
            ],
            help="Pulse Check output generates a Word document using the pulse-check template."
        )

        st.markdown("---")
        
        # Generate button
        if st.button("🚀 Generate Report", use_container_width=True):
            with st.spinner("Generating your report..."):
                if report_style == "Pulse Check (Word .docx)":
                    template_path = os.path.join(os.path.dirname(__file__), 'pulse_check_template.docx')
                    sections = _pulse_check_sections(em_filtered, sms_filtered, sl_filtered)
                    docx_output = build_pulse_check_docx(
                        template_path=template_path,
                        title_line=f"{campaign_name} - Engagement Pulse Check",
                        sections=sections
                    )
                    file_suffix = "pulse_check"
                    download_ext = "docx"
                    download_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                else:
                    # Load template
                    template_path = os.path.join(os.path.dirname(__file__), 'template.pptx')
                    with open(template_path, 'rb') as f:
                        template_bytes = f.read()
                    
                    # Generate PPT
                    ppt_output = generate_campaign_ppt(
                        em_filtered, 
                        sms_filtered, 
                        sl_filtered, 
                        campaign_name, 
                        template_bytes
                    )
                    docx_output = ppt_output
                    file_suffix = "report"
                    download_ext = "pptx"
                    download_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                
                st.markdown('<div class="success-box">✅ Report Generated Successfully!</div>', unsafe_allow_html=True)
                
                # Download button
                safe_filename = campaign_name.replace("'", "").replace(" ", "_")
                st.download_button(
                    label="📥 Download Report",
                    data=docx_output,
                    file_name=f"{safe_filename}_{file_suffix}.{download_ext}",
                    mime=download_mime,
                    use_container_width=True
                )
    
    except Exception as e:
        st.error(f"❌ Error: {str(e)}")
        st.info("Please make sure your Excel file has the correct tabs: EM, SMS, SLs")
        import traceback
        st.code(traceback.format_exc())

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #888; font-size: 0.9rem;">
    Made with ❤️ using Streamlit | T-Mobile Marketing Analytics
</div>
""", unsafe_allow_html=True)
