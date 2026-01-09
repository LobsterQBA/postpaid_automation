"""
Generates a PowerPoint report using template.pptx and functions from functions.py
"""

import functions as DS
from pptx import Presentation
import pandas as pd


def main():
    # ========== 1. SETUP & LOAD ==========
    pptx_file = "template.pptx"
    prs = Presentation(pptx_file)
    
    # ========== 2. MAP - Get slide positions ==========
    navigation_screen = DS.create_navigation_screen(prs)
    
    # ========== 3. SLIDE 1 - Title Slide ==========
    prs = DS.duplicate_slide(prs, navigation_screen['titleslide'], verbose=False)
    title_slide = prs.slides[-1]
    DS.find_replace_text(title_slide, 'PLACE_TEXT_TITLE', "Leo's Automated Weekly Report", verbose=False)
    
    # ========== 4. SLIDE 2 - Performance Overview ==========
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    data_slide = prs.slides[-1]
    DS.find_replace_text(data_slide, 'PLACE_TEXT_TITLE', 'Performance Overview', verbose=False)
    
    # Create mock DataFrame with English data
    mock_data = pd.DataFrame({
        'Touch': ['Launch', 'Preorder', 'Announce', 'Reminder'],
        'Audience': ['Growth_AAL', 'Churn_Upgrade', 'New_Users', 'Engaged_Users'],
        'Creative': ['Spring Campaign A', 'Early Access Promo', 'Product Reveal', 'Follow-up Message'],
        'Volume': ['1.2M', '850K', '2.5M', '500K'],
        'CTR': ['2.5%', '3.8%', '1.9%', '4.2%']
    })
    DS.add_data_table_new(data_slide, 'Touch', mock_data)
    
    # ========== 5. SECTION: Subject Line Test ==========
    # Section Title
    prs = DS.duplicate_slide(prs, navigation_screen['subtitleslide'], verbose=False)
    section_slide = prs.slides[-1]
    DS.find_replace_text(section_slide, 'PLACE_TEXT_TITLE', 'Subject Line Test Results', verbose=False)
    
    # Subject Line Test Data Slide
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    sl_slide = prs.slides[-1]
    DS.find_replace_text(sl_slide, 'PLACE_TEXT_TITLE', 'G3 Subject Line Results', verbose=False)
    
    sl_data = pd.DataFrame({
        'Touch': ['Apple Winner', 'Apple Loser', 'Google Winner', 'Google Loser'],
        'Audience': ['<FN>, your iPhone 17 is On Us', '<FN>, Get iPhone 17 On Us', '<FN>, your Pixel 10 is On Us', '<FN>, Get Pixel 10 On Us'],
        'Creative': ['Direct messaging', 'Call to action', 'Direct messaging', 'Call to action'],
        'Volume': ['538K', '538K', '17.7K', '17.7K'],
        'CTR': ['1.92%', '0.96%', '1.85%', '1.02%']
    })
    DS.add_data_table_new(sl_slide, 'Touch', sl_data)
    
    # ========== 6. SECTION: SMS Performance ==========
    # Section Title
    prs = DS.duplicate_slide(prs, navigation_screen['subtitleslide'], verbose=False)
    section_slide = prs.slides[-1]
    DS.find_replace_text(section_slide, 'PLACE_TEXT_TITLE', 'SMS/RBM Performance', verbose=False)
    
    # SMS Performance Data Slide
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    sms_slide = prs.slides[-1]
    DS.find_replace_text(sms_slide, 'PLACE_TEXT_TITLE', 'T1: SMS/RBM Performance', verbose=False)
    
    sms_data = pd.DataFrame({
        'Touch': ['Apple G1', 'Apple G2', 'Apple G3', 'Samsung G1'],
        'Audience': ['Tech Savvy', 'At-Risk', 'Holiday Gifting', 'Tech Savvy'],
        'Creative': ['iPhone 17 On Us - Plan Value', 'iPhone 17 - Thanks for being with us', 'iPhone 17 - Stay Connected', 'Galaxy S25 On Us - Samsung Tech'],
        'Volume': ['2.77%', '3.82%', '1.68%', '12.90%'],
        'CTR': ['6.60%', '14.5%', '3.97%', '2.77%']
    })
    DS.add_data_table_new(sms_slide, 'Touch', sms_data)
    
    # ========== 7. SECTION: Creative Engagement ==========
    # Section Title
    prs = DS.duplicate_slide(prs, navigation_screen['subtitleslide'], verbose=False)
    section_slide = prs.slides[-1]
    DS.find_replace_text(section_slide, 'PLACE_TEXT_TITLE', 'Creative Engagement Analysis', verbose=False)
    
    # Creative Engagement Email Slide
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    email_slide = prs.slides[-1]
    DS.find_replace_text(email_slide, 'PLACE_TEXT_TITLE', 'Creative Engagement: Email', verbose=False)
    
    email_data = pd.DataFrame({
        'Touch': ['Hero/1st', '2nd', '3rd', '4th'],
        'Audience': ['Get them now - Holiday Deals', 'Get yours now - Phone Offer', 'Get yours - Watch Offer', 'Device Accessory Icons'],
        'Creative': ['Holiday campaign CTA', 'Phone upgrade CTA', 'Watch promo CTA', 'Accessory showcase'],
        'Volume': ['0.12%', '3.10%', '0.24%', '0.07%'],
        'CTR': ['0.48%', '1.17%', '0.30%', '0.09%']
    })
    DS.add_data_table_new(email_slide, 'Touch', email_data)
    
    # Creative Engagement Heatmap Slide
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    heatmap_slide = prs.slides[-1]
    DS.find_replace_text(heatmap_slide, 'PLACE_TEXT_TITLE', 'Creative Engagement: Heatmap', verbose=False)
    
    heatmap_data = pd.DataFrame({
        'Touch': ['Module 1', 'Module 2', 'Module 3', 'Module 4'],
        'Audience': ['Hero Banner', 'Product Grid', 'CTA Button', 'Footer Links'],
        'Creative': ['High engagement zone', 'Medium engagement', 'Primary click area', 'Low engagement'],
        'Volume': ['45%', '28%', '22%', '5%'],
        'CTR': ['12.5%', '8.3%', '15.2%', '2.1%']
    })
    DS.add_data_table_new(heatmap_slide, 'Touch', heatmap_data)
    
    # ========== 8. CLEANUP - Delete original template slides ==========
    indices_to_delete = list(navigation_screen.values())
    for i in sorted(indices_to_delete, reverse=True):
        rId = prs.slides._sldIdLst[i]
        prs.slides._sldIdLst.remove(rId)
        print(f"Deleted template slide at index {i}")
    
    # ========== 9. SAVE ==========
    output_file = 'output_leo_report.pptx'
    prs.save(output_file)
    print(f"\nReport saved to: {output_file}")
    print("Total slides generated: ", len(prs.slides))


if __name__ == "__main__":
    main()
