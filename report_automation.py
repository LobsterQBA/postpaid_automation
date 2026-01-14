"""
Automated Report Generator for EM/SMS/SL Testing
Reads data from Template Excel and generates PowerPoint report
"""

import functions as DS
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import os


def load_template_data(filepath="data_template.xlsx"):
    """Load all sheets from template Excel file"""
    em_data = pd.read_excel(filepath, sheet_name='EM')
    sms_data = pd.read_excel(filepath, sheet_name='SMS')
    sl_data = pd.read_excel(filepath, sheet_name='SL Testing')
    return em_data, sms_data, sl_data


def format_number(x):
    """Format numbers with comma separators"""
    if pd.isna(x) or x == "":
        return ""
    try:
        x = int(float(x))
        return f"{x:,}"
    except:
        return str(x)


def format_percent(x):
    """Format as percentage"""
    if pd.isna(x) or x == "":
        return ""
    try:
        x = float(x)
        if x < 1:
            return f"{x*100:.2f}%"
        else:
            return f"{x:.2f}%"
    except:
        return str(x)


def update_table_headers(slide, new_headers):
    """Update table header row with new column names"""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for c, header in enumerate(new_headers):
                if c < len(table.columns):
                    cell = table.cell(0, c)
                    cell.text_frame.paragraphs[0].clear()
                    run = cell.text_frame.paragraphs[0].add_run()
                    run.text = header
            return


def add_data_to_table(slide, data, font_size=9):
    """Add data to table, dynamically adding rows if needed"""
    from pptx.util import Pt
    from copy import deepcopy
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            cols = len(table.columns)
            current_rows = len(table.rows)
            needed_rows = data.shape[0] + 1  # +1 for header
            
            # Add rows if needed
            while len(table.rows) < needed_rows:
                # Add a new row by copying the last row's XML
                tbl = table._tbl
                tr = tbl.tr_lst[-1]
                new_tr = deepcopy(tr)
                tbl.append(new_tr)
            
            # Update headers (row 0)
            for c, col_name in enumerate(data.columns):
                if c < cols:
                    cell = table.cell(0, c)
                    p = cell.text_frame.paragraphs[0]
                    p.clear()
                    run = p.add_run()
                    run.text = col_name
                    run.font.size = Pt(font_size + 1)
                    run.font.bold = True
            
            # Update data rows (row 1+)
            for r in range(data.shape[0]):
                for c in range(min(data.shape[1], cols)):
                    cell = table.cell(r + 1, c)
                    p = cell.text_frame.paragraphs[0]
                    text = str(data.iloc[r, c]) if pd.notna(data.iloc[r, c]) else ""
                    p.clear()
                    run = p.add_run()
                    run.text = text
                    run.font.size = Pt(font_size)
            
            # Clear unused columns
            for c in range(data.shape[1], cols):
                for r in range(len(table.rows)):
                    cell = table.cell(r, c)
                    cell.text_frame.paragraphs[0].clear()
            
            return


def process_em_data(em_data):
    """Process EM data for PPT table"""
    df = em_data.copy()
    
    # Calculate Open Rate and Click Rate
    df['Open Rate'] = df['Unique Opens'] / df['Deliveries']
    df['Click Rate'] = df['Unique Clicks'] / df['Deliveries']
    
    # Select columns: Group, Offer Type, Deliveries, Opens, Clicks, Open Rate, Click Rate, Creative Message (last)
    table_df = df[['Group', 'Offer type', 'Deliveries', 'Unique Opens', 'Unique Clicks', 'Open Rate', 'Click Rate', 'Subject Line']].copy()
    table_df.columns = ['Group', 'Offer Type', 'Deliveries', 'Opens', 'Clicks', 'Open Rate', 'Click Rate', 'Creative Message']
    
    # Format numbers (full numbers with commas)
    table_df['Deliveries'] = table_df['Deliveries'].apply(format_number)
    table_df['Opens'] = table_df['Opens'].apply(format_number)
    table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    table_df['Open Rate'] = table_df['Open Rate'].apply(format_percent)
    table_df['Click Rate'] = table_df['Click Rate'].apply(format_percent)
    
    return table_df


def process_sms_data(sms_data):
    """Process SMS data for PPT table"""
    df = sms_data.copy()
    
    # Select columns: Touch, Channel, Group, Sent, Clicks, CTR, Creative Message (last)
    table_df = df[['Touch', 'Channel', 'Group', 'Sent (Est.)', 'Clicks', 'CTR', 'Creative Message']].copy()
    table_df.columns = ['Touch', 'Channel', 'Group', 'Sent', 'Clicks', 'CTR', 'Creative Message']
    
    # Format numbers (full numbers with commas)
    table_df['Sent'] = table_df['Sent'].apply(format_number)
    table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    table_df['CTR'] = table_df['CTR'].apply(format_percent)
    
    return table_df


def process_sl_testing_data(sl_data):
    """Process SL Testing data and determine winner"""
    # Filter out Grand Total row
    df = sl_data[sl_data['Row Labels'] != 'Grand Total'].copy()
    
    # Rename columns
    df = df.rename(columns={
        'Row Labels': 'Variant',
        'Sum of Unique People Clicks [v8]': 'Clicks',
        'Sum of People Opens [v8]': 'Opens',
        'Sum of Unique People Delivered [v8]': 'Delivered',
        'Open Rate': 'Open Rate',
        'Click Rate': 'Click Rate',
        'Subject Line': 'Subject Line'
    })
    
    # Determine winner based on Open Rate
    winner_idx = df['Open Rate'].idxmax()
    loser_idx = df['Open Rate'].idxmin()
    
    winner_row = df.loc[winner_idx]
    loser_row = df.loc[loser_idx]
    
    # Calculate percentage difference: (winner - loser) / loser
    open_rate_diff = (winner_row['Open Rate'] - loser_row['Open Rate']) / loser_row['Open Rate']
    ctr_diff = (winner_row['Click Rate'] - loser_row['Click Rate']) / loser_row['Click Rate']
    
    winner_info = {
        'variant': winner_row['Variant'],
        'subject_line': winner_row['Subject Line'],
        'open_rate': winner_row['Open Rate'],
        'ctr': winner_row['Click Rate'],
        'open_rate_diff': open_rate_diff,
        'ctr_diff': ctr_diff
    }
    
    # Prepare table data: Variant, Delivered, Opens, Open Rate, Clicks, Click Rate, Subject Line (last)
    table_df = df[['Variant', 'Delivered', 'Opens', 'Open Rate', 'Clicks', 'Click Rate', 'Subject Line']].copy()
    
    # Format numbers (full numbers with commas)
    table_df['Delivered'] = table_df['Delivered'].apply(format_number)
    table_df['Opens'] = table_df['Opens'].apply(format_number)
    table_df['Clicks'] = table_df['Clicks'].apply(format_number)
    table_df['Open Rate'] = table_df['Open Rate'].apply(format_percent)
    table_df['Click Rate'] = table_df['Click Rate'].apply(format_percent)
    
    return table_df, winner_info


def main():
    # ========== 1. SETUP & LOAD TEMPLATE ==========
    pptx_file = "template.pptx"
    prs = Presentation(pptx_file)
    navigation_screen = DS.create_navigation_screen(prs)
    
    # ========== 2. LOAD DATA ==========
    print("\n--- Loading Data ---")
    em_data, sms_data, sl_data = load_template_data("data_template.xlsx")
    print(f"EM data: {len(em_data)} rows")
    print(f"SMS data: {len(sms_data)} rows")
    print(f"SL Testing data: {len(sl_data)} rows")
    
    # ========== 3. TITLE SLIDE ==========
    prs = DS.duplicate_slide(prs, navigation_screen['titleslide'], verbose=False)
    title_slide = prs.slides[-1]
    DS.find_replace_text(title_slide, 'PLACE_TEXT_TITLE', "Weekly Report", verbose=False)
    
    # ========== 4. EMAIL PERFORMANCE SECTION ==========
    prs = DS.duplicate_slide(prs, navigation_screen['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Performance', verbose=False)
    
    # EM Data slide
    em_table = process_em_data(em_data)
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    em_slide = prs.slides[-1]
    DS.find_replace_text(em_slide, 'PLACE_TEXT_TITLE', 'Email Performance Overview', verbose=False)
    
    # Add EM data (all rows)
    add_data_to_table(em_slide, em_table)
    print(f"EM table: {em_table.shape[0]} rows x {em_table.shape[1]} cols")
    
    # ========== 5. SMS/RBM PERFORMANCE SECTION ==========
    prs = DS.duplicate_slide(prs, navigation_screen['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'SMS/RBM Performance', verbose=False)
    
    # SMS Data slide
    sms_table = process_sms_data(sms_data)
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    sms_slide = prs.slides[-1]
    DS.find_replace_text(sms_slide, 'PLACE_TEXT_TITLE', 'SMS/RBM Performance Overview', verbose=False)
    
    # Add SMS data (all rows)
    add_data_to_table(sms_slide, sms_table)
    print(f"SMS table: {sms_table.shape[0]} rows x {sms_table.shape[1]} cols")
    
    # ========== 6. SUBJECT LINE TESTING SECTION ==========
    # SL Testing Data slide (no subtitle, no winner box)
    sl_table, winner_info = process_sl_testing_data(sl_data)
    prs = DS.duplicate_slide(prs, navigation_screen['smsdataslide'], verbose=False)
    sl_slide = prs.slides[-1]
    DS.find_replace_text(sl_slide, 'PLACE_TEXT_TITLE', 'Subject Line Testing Results', verbose=False)
    
    # Add SL data (all rows)
    add_data_to_table(sl_slide, sl_table)
    print(f"SL table: {sl_table.shape[0]} rows x {sl_table.shape[1]} cols")
    
    # ========== 7. CLEANUP ==========
    for i in sorted(navigation_screen.values(), reverse=True):
        rId = prs.slides._sldIdLst[i]
        prs.slides._sldIdLst.remove(rId)
        print(f"Deleted template slide at index {i}")
    
    # ========== 8. SAVE ==========
    output_file = 'output_report.pptx'
    prs.save(output_file)
    print(f"\nReport saved to: {output_file}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
