"""
Auto Load Demo — Generic Campaign Report Generator
====================================================
One-click pipeline: Load cleaned data (EM + RCM SMS + SL Testing) → Generate PPT or Pulse Check Doc

Supports any campaign with the standard data format:
  - EM / RCM SMS / SL Testing (derived or explicit)
  - Standard SMS or RCM dual-channel (auto-detected)

Output formats:
  - Campaign Report (PPT)
  - Pulse Check (Word .docx)

Usage:
    streamlit run auto_load_demo.py
"""

import streamlit as st
import pandas as pd
import numpy as np
import io
import os
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from copy import deepcopy
from pulse_check_docx import build_pulse_check_docx

# ========== PAGE CONFIG ==========
st.set_page_config(page_title="Auto Load Demo", page_icon="🚀", layout="wide")

# ========== T-Mobile Magenta Theme ==========
st.markdown("""
<style>
    :root { --magenta: #E20074; --magenta-dark: #B8005C; }
    .main-header {
        background: linear-gradient(135deg, #E20074 0%, #B8005C 100%);
        padding: 2rem; border-radius: 10px; text-align: center; margin-bottom: 2rem;
    }
    .main-header h1 { color: white; margin: 0; font-size: 2.2rem; }
    .main-header p { color: rgba(255,255,255,0.9); margin: 0.5rem 0 0 0; font-size: 1rem; }
    .stButton > button {
        background: linear-gradient(135deg, #E20074 0%, #B8005C 100%);
        color: white; border: none; padding: 0.75rem 2rem;
        font-size: 1.1rem; font-weight: 600; border-radius: 8px; width: 100%;
    }
    .stButton > button:hover {
        background: linear-gradient(135deg, #FF1493 0%, #E20074 100%);
        box-shadow: 0 4px 15px rgba(226, 0, 116, 0.4);
    }
    .info-box {
        background: #f8f9fa; border-left: 4px solid #E20074;
        padding: 1rem; border-radius: 0 8px 8px 0; margin: 1rem 0;
    }
    .success-box {
        background: linear-gradient(135deg, #00C853 0%, #00E676 100%);
        padding: 1rem; border-radius: 8px; text-align: center; color: white; font-weight: 600;
    }
</style>
""", unsafe_allow_html=True)

# ========== HEADER ==========
st.markdown("""
<div class="main-header">
    <h1>🚀 Auto Load Demo</h1>
    <p>One-click: Load Cleaned Data → Generate Campaign Report (PPT) or Pulse Check (Word)</p>
</div>
""", unsafe_allow_html=True)


# =====================================================================
# HELPER FUNCTIONS
# =====================================================================

def format_number(x):
    """Format a number with comma separators. Handles NaN/empty gracefully."""
    if pd.isna(x) or x == "":
        return ""
    try:
        return f"{int(float(x)):,}"
    except (ValueError, TypeError):
        return str(x)


def format_percent(x):
    """
    Format a value as percentage string.
    Smart detection: values < 1 are treated as ratios (0.29 → 29.00%),
    values >= 1 are treated as already-percentage (29.07 → 29.07%).
    Edge case: exactly 1.0 is treated as 100%.
    """
    if pd.isna(x) or x == "":
        return ""
    try:
        x = float(x)
        if x == 0:
            return "0.00%"
        # Heuristic: if absolute value <= 1, it's a ratio → multiply by 100
        # This covers typical marketing rates (OR < 100%, CTR < 100%)
        if abs(x) <= 1:
            return f"{x * 100:.2f}%"
        else:
            return f"{x:.2f}%"
    except (ValueError, TypeError):
        return str(x)


def truncate_text(text, max_length=50):
    """Truncate text to max_length with ellipsis."""
    if pd.isna(text):
        return ""
    text = str(text).strip()
    return text[:max_length - 3] + "..." if len(text) > max_length else text


def is_column_empty(df, col_name):
    """Check if a column is entirely empty/NaN."""
    if col_name not in df.columns:
        return True
    return df[col_name].isna().all() or (df[col_name].astype(str).str.strip() == '').all()


def _safe_div(a, b):
    """Safe division, returns 0 when dividing by zero."""
    try:
        return a / b if b else 0
    except (TypeError, ZeroDivisionError):
        return 0


def _find_column(df, candidates, default=None):
    """Find the first matching column name from a list of candidates."""
    for c in candidates:
        if c in df.columns:
            return c
    return default


def _ensure_touch_column(df):
    """Ensure DataFrame has a 'Touch' column; create default if missing."""
    if df is None:
        return df
    if 'Touch' not in df.columns:
        df = df.copy()
        df['Touch'] = 'Touch 1'
    return df


# =====================================================================
# DATA VALIDATION
# =====================================================================

def validate_data(em_data, sms_data, sms_format):
    """
    Validate loaded data and return a list of warnings (non-fatal)
    and a list of errors (fatal).
    """
    warnings = []
    errors = []

    if em_data is None and sms_data is None:
        errors.append("No EM or SMS/RCM data found. Need at least one.")
        return warnings, errors

    # EM validation
    if em_data is not None:
        required_em = ['Deliveries', 'Unique Opens', 'Unique Clicks']
        missing = [c for c in required_em if c not in em_data.columns]
        if missing:
            errors.append(f"EM tab missing required columns: {missing}")
        if 'Touch' not in em_data.columns:
            warnings.append("EM tab has no 'Touch' column — defaulting all rows to 'Touch 1'")

    # SMS/RCM validation
    if sms_data is not None:
        if sms_format == "rcm":
            # RCM format needs at least delivery columns
            has_sms_del = 'SMS Deliveries' in sms_data.columns
            has_rbm_del = 'RBM Deliveries' in sms_data.columns
            if not has_sms_del and not has_rbm_del:
                errors.append("RCM SMS tab missing delivery columns (SMS Deliveries / RBM Deliveries)")
        else:
            if 'Deliveries' not in sms_data.columns:
                errors.append("SMS tab missing 'Deliveries' column")
            clicks_col = _find_column(sms_data,
                                      ['PBI Unique Clicks', 'Branch Unique Clicks', 'Unique Clicks'])
            if clicks_col is None:
                warnings.append("SMS tab has no recognizable clicks column — CTR will show 0%")

        if 'Touch' not in sms_data.columns:
            warnings.append("SMS/RCM tab has no 'Touch' column — defaulting all rows to 'Touch 1'")

    return warnings, errors


# =====================================================================
# SMART DATA LOADER — auto-detect sheet format
# =====================================================================

def detect_and_load(file_or_path):
    """
    Auto-detect sheet names and return a dict with standardised keys:
      em_data, sms_data, sl_data (may be None)
    Supports:
      Standard format:  EM / SMS / SLs
      RCM format:       EM / RCM SMS  (SL derived from EM if no SLs sheet)
    """
    xls = pd.ExcelFile(file_or_path)
    sheets = xls.sheet_names

    result = {"em_data": None, "sms_data": None, "sl_data": None,
              "sms_format": "standard", "sheets_found": sheets}

    # --- EM ---
    em_sheet = _find_column(pd.DataFrame(columns=sheets),
                            []) # dummy — use direct check below
    for name in ['EM', 'Email', 'EM Data', 'Email Data']:
        if name in sheets:
            result["em_data"] = pd.read_excel(xls, sheet_name=name)
            break

    # --- SMS / RCM SMS ---
    for name in ['RCM SMS', 'SMS', 'RCM', 'SMS Data']:
        if name in sheets:
            result["sms_data"] = pd.read_excel(xls, sheet_name=name)
            # Auto-detect RCM format: by sheet name OR by presence of RCM-specific columns
            if "RCM" in name:
                result["sms_format"] = "rcm"
            elif result["sms_data"] is not None:
                rcm_cols = ['SMS Deliveries', 'RBM Deliveries', 'RBM InfoBip Clicks']
                if any(c in result["sms_data"].columns for c in rcm_cols):
                    result["sms_format"] = "rcm"
            break

    # --- SLs (explicit sheet) ---
    for name in ['SLs', 'SL Testing', 'Subject Lines', 'SL', 'SL Test']:
        if name in sheets:
            result["sl_data"] = pd.read_excel(xls, sheet_name=name)
            break

    # Ensure Touch column exists
    result["em_data"] = _ensure_touch_column(result["em_data"])
    result["sms_data"] = _ensure_touch_column(result["sms_data"])

    return result


# =====================================================================
# DATA PROCESSING — Section 1: Campaign Overview (Engagement Summary)
# =====================================================================

def process_engagement_summary(em_data, sms_data, sms_format="standard"):
    """
    Build the top-level Campaign Overview engagement summary table.
    Columns adapt based on available data:
      Touch | Email OR | Email CTR | Email Delivery | [RCM CTR | SMS CTR] or [SMS CTR]
    """
    touches = []
    if em_data is not None:
        touches.extend(em_data['Touch'].unique().tolist())
    if sms_data is not None:
        touches.extend(sms_data['Touch'].unique().tolist())
    touches = list(dict.fromkeys(touches))  # deduplicate, preserve order

    rows = []
    for touch in touches:
        row = {'Touch': touch}

        # EM metrics
        if em_data is not None:
            em_touch = em_data[em_data['Touch'] == touch]
            if len(em_touch) > 0:
                total_del = em_touch['Deliveries'].sum()
                row['Email OR'] = format_percent(_safe_div(em_touch['Unique Opens'].sum(), total_del))
                row['Email CTR'] = format_percent(_safe_div(em_touch['Unique Clicks'].sum(), total_del))
                row['Email Delivery'] = format_number(total_del)
            else:
                row['Email OR'] = ''
                row['Email CTR'] = ''
                row['Email Delivery'] = ''
        else:
            row['Email OR'] = ''
            row['Email CTR'] = ''
            row['Email Delivery'] = ''

        # SMS / RCM metrics
        if sms_data is not None:
            sms_touch = sms_data[sms_data['Touch'] == touch]
            if len(sms_touch) > 0:
                if sms_format == "rcm":
                    sms_del = sms_touch['SMS Deliveries'].sum() if 'SMS Deliveries' in sms_touch.columns else 0
                    rbm_del = sms_touch['RBM Deliveries'].sum() if 'RBM Deliveries' in sms_touch.columns else 0
                    sms_click_col = _find_column(sms_touch, ['SMS InfoBip Clicks', 'SMS Clicks'])
                    rbm_click_col = _find_column(sms_touch, ['RBM InfoBip Clicks', 'RBM Clicks'])
                    sms_clicks = sms_touch[sms_click_col].sum() if sms_click_col else 0
                    rbm_clicks = sms_touch[rbm_click_col].sum() if rbm_click_col else 0
                    row['RCM CTR'] = format_percent(_safe_div(rbm_clicks, rbm_del))
                    row['SMS CTR'] = format_percent(_safe_div(sms_clicks, sms_del))
                else:
                    total_del = sms_touch['Deliveries'].sum() if 'Deliveries' in sms_touch.columns else 0
                    clicks_col = _find_column(sms_touch,
                                              ['PBI Unique Clicks', 'Branch Unique Clicks', 'Unique Clicks'])
                    total_clicks = sms_touch[clicks_col].sum() if clicks_col else 0
                    row['SMS CTR'] = format_percent(_safe_div(total_clicks, total_del))
            else:
                if sms_format == "rcm":
                    row['RCM CTR'] = ''
                    row['SMS CTR'] = ''
                else:
                    row['SMS CTR'] = ''
        else:
            if sms_format == "rcm":
                row['RCM CTR'] = ''
                row['SMS CTR'] = ''
            else:
                row['SMS CTR'] = ''

        rows.append(row)

    if not rows:
        return pd.DataFrame()

    # Build DataFrame — only include columns that have at least one non-empty value
    if sms_format == "rcm":
        col_order = ['Touch', 'Email OR', 'Email CTR', 'Email Delivery', 'RCM CTR', 'SMS CTR']
    else:
        col_order = ['Touch', 'Email OR', 'Email CTR', 'Email Delivery', 'SMS CTR']

    df = pd.DataFrame(rows)

    # Drop columns that are entirely empty (e.g., no EM data → drop Email columns)
    for col in list(col_order):
        if col in df.columns and col != 'Touch':
            if df[col].replace('', np.nan).isna().all():
                col_order.remove(col)

    col_order = [c for c in col_order if c in df.columns]
    return df[col_order]


# =====================================================================
# DATA PROCESSING — Section 2: Email High-Level Results
# =====================================================================

def process_email_summary_by_touch(em_data):
    """Email summary aggregated by Touch."""
    if em_data is None or len(em_data) == 0:
        return pd.DataFrame(columns=['Touch', 'Deliveries', 'Open Rate', 'CTR'])
    summary = em_data.groupby('Touch').agg(
        {'Deliveries': 'sum', 'Unique Opens': 'sum', 'Unique Clicks': 'sum'}
    ).reset_index()
    summary['Open Rate'] = summary['Unique Opens'] / summary['Deliveries']
    summary['CTR'] = summary['Unique Clicks'] / summary['Deliveries']
    return pd.DataFrame({
        'Touch': summary['Touch'],
        'Deliveries': summary['Deliveries'].apply(format_number),
        'Open Rate': summary['Open Rate'].apply(format_percent),
        'CTR': summary['CTR'].apply(format_percent)
    })


def process_em_detail(em_data):
    """
    Full email detail table. Columns are dynamic based on what's in the data:
    Touch | OS | Cohort | Audience Details 1/2/3 | Deliveries | Opens | Clicks | OR | CTR
    """
    if em_data is None or len(em_data) == 0:
        return pd.DataFrame()

    df = em_data.copy()
    df['Open Rate'] = df['Unique Opens'] / df['Deliveries']
    df['Click Rate'] = df['Unique Clicks'] / df['Deliveries']

    # Build column list dynamically — only include non-empty columns
    display_cols = []
    for col in ['Touch', 'OS', 'Cohort']:
        if col in df.columns and not is_column_empty(df, col):
            display_cols.append(col)
    for col in ['Audience Details 1', 'Audience Details 2', 'Audience Details 3']:
        if col in df.columns and not is_column_empty(df, col):
            display_cols.append(col)
    display_cols.extend(['Deliveries', 'Unique Opens', 'Unique Clicks', 'Open Rate', 'Click Rate'])
    display_cols = [c for c in display_cols if c in df.columns]

    table_df = df[display_cols].copy()
    for c in ['Deliveries', 'Unique Opens', 'Unique Clicks']:
        if c in table_df.columns:
            table_df[c] = table_df[c].apply(format_number)
    table_df['Open Rate'] = table_df['Open Rate'].apply(format_percent)
    table_df['Click Rate'] = table_df['Click Rate'].apply(format_percent)
    return table_df


# =====================================================================
# DATA PROCESSING — Section 3: Subject Line Testing
# =====================================================================

def _strip_variant_suffix(label):
    """
    Strip trailing SLA/SLB/SLC/etc. from a Delivery Label to get the test group key.
    e.g. '4GPiOSnolangprefEMSLA' → '4GPiOSnolangprefEM'
    """
    import re
    return re.sub(r'SL[A-Z]+$', '', str(label).strip())


def _build_slide_label(touch, sub, audience_col):
    """
    Build a human-readable slide label from Touch + audience fields.
    e.g. 'Touch 1 — Growth / Prime'
    """
    parts = []
    for col in ['Cohort', 'Audience Details 1', 'Audience Details 2', 'Audience Details 3']:
        if col in sub.columns and not is_column_empty(sub, col):
            val = sub[col].dropna().iloc[0]
            if str(val).strip():
                parts.append(str(val).strip())
    audience_desc = ' / '.join(parts) if parts else (
        sub[audience_col].dropna().iloc[0] if audience_col and audience_col in sub.columns else 'All'
    )
    return f"{touch} — {audience_desc}"


def process_sl_testing(sl_data):
    """
    Build SL testing tables grouped by Touch + Delivery Label prefix.
    Each unique (Touch, prefix) = 1 slide, with one row per SL variant (SLA/SLB/...).
    Falls back to (Touch, Audience) grouping if no Delivery Label column exists.
    Returns a list of (label, DataFrame) tuples.
    Each DataFrame: Variant | Subject Line | Delivery | Open Rate | Click Rate | OR Lift | CTR Lift
    """
    if sl_data is None or len(sl_data) == 0:
        return []

    df = sl_data.copy()

    # Determine audience grouping column (for slide labels)
    if 'Audience' in df.columns:
        audience_col = 'Audience'
    elif 'Cohort' in df.columns:
        audience_col = 'Cohort'
    else:
        audience_col = None

    has_touch = 'Touch' in df.columns and not is_column_empty(df, 'Touch')
    has_delivery_label = 'Delivery Label' in df.columns and not is_column_empty(df, 'Delivery Label')

    # Normalise variant column name
    # Note: 'Delivery Label' is NOT a variant — it's the test group identifier
    variant_col = _find_column(df, ['SL Testing Variant', 'Variant', 'Row Labels',
                                     'Treatment', 'Test Group'])
    if variant_col and variant_col != 'Variant':
        df['Variant'] = df[variant_col]

    # Normalise rate column names
    if 'OR' in df.columns and 'Open Rate' not in df.columns:
        df['Open Rate'] = df['OR']
    if 'CTR' in df.columns and 'Click Rate' not in df.columns:
        df['Click Rate'] = df['CTR']
    if 'Delivered' in df.columns and 'Deliveries' not in df.columns:
        df['Deliveries'] = df['Delivered']

    results = []

    if has_touch and has_delivery_label:
        # PRIMARY path: group by Touch, then by Delivery Label prefix (strip SLA/SLB suffix)
        # This correctly handles cases where multiple test groups share the same Audience
        df['_test_group'] = df['Delivery Label'].apply(_strip_variant_suffix)
        for touch in df['Touch'].unique():
            touch_df = df[df['Touch'] == touch]
            for prefix in touch_df['_test_group'].unique():
                sub = touch_df[touch_df['_test_group'] == prefix].copy()
                table = _build_sl_table_aggregated(sub)
                if table is not None and len(table) > 0:
                    label = _build_slide_label(touch, sub, audience_col)
                    results.append((label, table))

    elif has_touch and audience_col:
        # FALLBACK: no Delivery Label — group by (Touch, Audience)
        for touch in df['Touch'].unique():
            touch_df = df[df['Touch'] == touch]
            for audience in touch_df[audience_col].unique():
                sub = touch_df[touch_df[audience_col] == audience].copy()
                table = _build_sl_table_aggregated(sub)
                if table is not None and len(table) > 0:
                    results.append((f"{touch} — {audience}", table))

    elif has_touch:
        # FALLBACK: Touch only
        for touch in df['Touch'].unique():
            sub = df[df['Touch'] == touch].copy()
            table = _build_sl_table_aggregated(sub)
            if table is not None and len(table) > 0:
                results.append((str(touch), table))

    elif audience_col:
        # FALLBACK: Audience only
        for audience in df[audience_col].unique():
            sub = df[df[audience_col] == audience].copy()
            table = _build_sl_table_aggregated(sub)
            if table is not None and len(table) > 0:
                results.append((str(audience), table))

    else:
        table = _build_sl_table_aggregated(df)
        if table is not None and len(table) > 0:
            results.append(("All", table))

    return results


def _build_sl_table_aggregated(df):
    """
    Build a single SL testing table, aggregating duplicate rows per variant.
    Output: exactly one row per variant with weighted-average rates.
    Columns: Variant | Subject Line | Delivery | Open Rate | Click Rate | OR Lift | CTR Lift
    Lift is calculated for both Open Rate and Click Rate independently.
    For each metric, the variant with the highest rate gets a lift value:
    Lift = (best - second_best) / second_best, shown on the best variant's row only.
    """
    if 'Variant' not in df.columns:
        return None

    # Ensure numeric columns
    for c in ['Deliveries', 'Unique Opens', 'Unique Clicks']:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors='coerce').fillna(0)

    # If raw opens/clicks available, aggregate from those (more accurate)
    has_raw = 'Unique Opens' in df.columns and 'Unique Clicks' in df.columns

    rows = []
    for variant in df['Variant'].unique():
        vdf = df[df['Variant'] == variant]
        row = {'Variant': variant}

        # Subject Line — take the first non-null
        if 'Subject Line' in vdf.columns:
            sl_vals = vdf['Subject Line'].dropna()
            row['Subject Line'] = truncate_text(sl_vals.iloc[0], 50) if len(sl_vals) > 0 else ''
        else:
            row['Subject Line'] = ''

        total_del = vdf['Deliveries'].sum() if 'Deliveries' in vdf.columns else 0
        row['Delivery'] = format_number(total_del)

        # Compute raw OR/CTR for winner/lift calculation
        if has_raw:
            total_opens = vdf['Unique Opens'].sum()
            total_clicks = vdf['Unique Clicks'].sum()
            raw_or = _safe_div(total_opens, total_del)
            raw_ctr = _safe_div(total_clicks, total_del)
        elif 'Open Rate' in vdf.columns:
            if total_del > 0:
                raw_or = (vdf['Open Rate'] * vdf['Deliveries']).sum() / total_del
            else:
                raw_or = vdf['Open Rate'].mean()
            if 'Click Rate' in vdf.columns:
                if total_del > 0:
                    raw_ctr = (vdf['Click Rate'] * vdf['Deliveries']).sum() / total_del
                else:
                    raw_ctr = vdf['Click Rate'].mean()
            else:
                raw_ctr = 0
        else:
            raw_or = 0
            raw_ctr = 0

        row['Open Rate'] = format_percent(raw_or)
        row['Click Rate'] = format_percent(raw_ctr)
        row['_raw_or'] = raw_or    # keep raw value for lift calc
        row['_raw_ctr'] = raw_ctr  # keep raw value for lift calc
        rows.append(row)

    if not rows:
        return None

    # Determine Lift for both Open Rate and Click Rate
    # For each metric, the variant with the highest rate gets a lift value
    # Lift = (best - second_best) / second_best
    if len(rows) >= 2:
        # OR Lift: find the variant with the highest raw OR
        best_or_idx = max(range(len(rows)), key=lambda i: rows[i]['_raw_or'])
        other_ors = [rows[i]['_raw_or'] for i in range(len(rows)) if i != best_or_idx]
        second_or = max(other_ors) if other_ors else 0

        # CTR Lift: find the variant with the highest raw CTR
        best_ctr_idx = max(range(len(rows)), key=lambda i: rows[i]['_raw_ctr'])
        other_ctrs = [rows[i]['_raw_ctr'] for i in range(len(rows)) if i != best_ctr_idx]
        second_ctr = max(other_ctrs) if other_ctrs else 0

        for i, row in enumerate(rows):
            # OR Lift
            if i == best_or_idx and second_or > 0:
                or_lift = (row['_raw_or'] - second_or) / second_or
                row['OR Lift'] = f"+{or_lift * 100:.1f}%"
            else:
                row['OR Lift'] = ''

            # CTR Lift
            if i == best_ctr_idx and second_ctr > 0:
                ctr_lift = (row['_raw_ctr'] - second_ctr) / second_ctr
                row['CTR Lift'] = f"+{ctr_lift * 100:.1f}%"
            else:
                row['CTR Lift'] = ''
    else:
        # Only one variant — no comparison possible
        for row in rows:
            row['OR Lift'] = ''
            row['CTR Lift'] = ''

    # Remove internal raw fields
    for row in rows:
        del row['_raw_or']
        del row['_raw_ctr']

    col_order = ['Variant', 'Subject Line', 'Delivery', 'Open Rate', 'Click Rate', 'OR Lift', 'CTR Lift']
    result = pd.DataFrame(rows)
    col_order = [c for c in col_order if c in result.columns]
    return result[col_order]


# =====================================================================
# DATA PROCESSING — Section 4: RCM / SMS Engagement
# =====================================================================

def process_rcm_detail(sms_data, sms_format="standard"):
    """
    RCM/SMS detail table by audience. Columns adapt to format:
      RCM format  → Audience | RBM Sent | SMS Sent | RBM CTR | SMS CTR | [Creative]
      Standard    → Touch | OS | Cohort | [Audience Details] | Deliveries | Clicks | CTR
    """
    if sms_data is None or len(sms_data) == 0:
        return pd.DataFrame()

    df = sms_data.copy()

    if sms_format == "rcm":
        # Build audience label dynamically from available detail columns
        parts = []
        for col in ['Audience Details 1', 'Audience Details 2', 'Audience Details 3']:
            if col in df.columns and not is_column_empty(df, col):
                parts.append(col)
        if parts:
            df['Audience'] = df[parts].apply(
                lambda r: '\n'.join([str(v) for v in r if pd.notna(v) and str(v).strip()]), axis=1
            )
        elif 'Cohort' in df.columns:
            df['Audience'] = df['Cohort']
        else:
            df['Audience'] = 'All'

        # Compute CTRs
        rbm_click_col = _find_column(df, ['RBM InfoBip Clicks', 'RBM Clicks'])
        sms_click_col = _find_column(df, ['SMS InfoBip Clicks', 'SMS Clicks'])
        df['RBM CTR'] = df.apply(lambda r: _safe_div(
            r.get(rbm_click_col, 0) if rbm_click_col else 0,
            r.get('RBM Deliveries', 0)), axis=1)
        df['SMS CTR'] = df.apply(lambda r: _safe_div(
            r.get(sms_click_col, 0) if sms_click_col else 0,
            r.get('SMS Deliveries', 0)), axis=1)

        display_cols = ['Audience']
        if 'RBM Deliveries' in df.columns:
            df['RBM Sent'] = df['RBM Deliveries']
            display_cols.append('RBM Sent')
        if 'SMS Deliveries' in df.columns:
            df['SMS Sent'] = df['SMS Deliveries']
            display_cols.append('SMS Sent')
        display_cols.extend(['RBM CTR', 'SMS CTR'])
        if 'Creative' in df.columns and not is_column_empty(df, 'Creative'):
            display_cols.append('Creative')

        table_df = df[display_cols].copy()
        for c in ['RBM Sent', 'SMS Sent']:
            if c in table_df.columns:
                table_df[c] = table_df[c].apply(format_number)
        table_df['RBM CTR'] = table_df['RBM CTR'].apply(format_percent)
        table_df['SMS CTR'] = table_df['SMS CTR'].apply(format_percent)
        return table_df

    else:
        # Standard SMS format
        clicks_col = _find_column(df, ['PBI Unique Clicks', 'Branch Unique Clicks', 'Unique Clicks'])
        ctr_col = _find_column(df, ['PBI CTR', 'PBI  CTR', 'CTR'])

        if clicks_col:
            df['Clicks'] = df[clicks_col]
        if ctr_col and 'Clicks' not in df.columns:
            df['CTR_display'] = df[ctr_col]

        display_cols = [c for c in ['Touch', 'OS', 'Cohort'] if c in df.columns and not is_column_empty(df, c)]
        for col in ['Audience Details 1', 'Audience Details 2', 'Audience Details 3']:
            if col in df.columns and not is_column_empty(df, col):
                display_cols.append(col)
        if 'Creative' in df.columns and not is_column_empty(df, 'Creative'):
            display_cols.append('Creative')

        if 'Clicks' in df.columns:
            display_cols.extend(['Deliveries', 'Clicks'])
            # Compute CTR from clicks/deliveries
            df['CTR'] = df.apply(lambda r: _safe_div(r.get('Clicks', 0), r.get('Deliveries', 0)), axis=1)
            display_cols.append('CTR')
        elif 'CTR_display' in df.columns:
            display_cols.extend(['Deliveries', 'CTR_display'])

        display_cols = [c for c in display_cols if c in df.columns]

        table_df = df[display_cols].copy()
        if 'Deliveries' in table_df.columns:
            table_df['Deliveries'] = table_df['Deliveries'].apply(format_number)
        if 'Clicks' in table_df.columns:
            table_df['Clicks'] = table_df['Clicks'].apply(format_number)
        if 'CTR' in table_df.columns:
            table_df['CTR'] = table_df['CTR'].apply(format_percent)
        if 'CTR_display' in table_df.columns:
            table_df['CTR_display'] = table_df['CTR_display'].apply(format_percent)
            table_df.rename(columns={'CTR_display': 'CTR'}, inplace=True)
        return table_df


# =====================================================================
# PPT GENERATION — Core Engine
# =====================================================================

def _import_functions():
    """Import the existing functions.py module."""
    import importlib.util
    func_path = os.path.join(os.path.dirname(__file__), 'functions.py')
    spec = importlib.util.spec_from_file_location("functions", func_path)
    DS = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(DS)
    return DS


def _clean_slide_placeholders(slide):
    """Remove INSIGHTS and HEADLINE placeholder text from duplicated data slides."""
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            text = shape.text.strip()
            if text in ('INSIGHTS', 'HEADLINE'):
                shape.text_frame.clear()


def add_data_to_table(slide, data, font_size=9):
    """Add data to table on slide, dynamically resizing rows and columns."""
    _clean_slide_placeholders(slide)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            needed_rows = data.shape[0] + 1  # +1 for header

            # Resize rows
            while len(table.rows) > needed_rows:
                tbl = table._tbl
                tbl.remove(tbl.tr_lst[-1])
            while len(table.rows) < needed_rows:
                tbl = table._tbl
                new_tr = deepcopy(tbl.tr_lst[-1])
                tbl.append(new_tr)

            # Resize columns — remove extras
            while len(table.columns) > data.shape[1]:
                tbl = table._tbl
                gridCol_lst = tbl.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gridCol')
                if gridCol_lst:
                    gridCol_lst[-1].getparent().remove(gridCol_lst[-1])
                    for tr in tbl.tr_lst:
                        tcs = tr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}tc')
                        if tcs:
                            tr.remove(tcs[-1])

            cols = len(table.columns)

            # Write headers
            for c, col_name in enumerate(data.columns):
                if c < cols:
                    cell = table.cell(0, c)
                    p = cell.text_frame.paragraphs[0]
                    p.clear()
                    run = p.add_run()
                    run.text = str(col_name)
                    run.font.size = Pt(font_size + 1)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)

            # Header row background
            for i in range(cols):
                cell = table.cell(0, i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(226, 0, 116)  # T-Mobile Magenta

            # Write data rows
            for r in range(data.shape[0]):
                for c in range(min(data.shape[1], cols)):
                    cell = table.cell(r + 1, c)
                    p = cell.text_frame.paragraphs[0]
                    val = data.iloc[r, c]
                    text = str(val) if pd.notna(val) else ""
                    p.clear()
                    run = p.add_run()
                    run.text = text
                    run.font.size = Pt(font_size)

            # Auto-fit column widths
            for col_idx in range(cols):
                max_len = 0
                for row_idx in range(len(table.rows)):
                    cell = table.cell(row_idx, col_idx)
                    max_len = max(max_len, len(str(cell.text)) if cell.text else 0)
                width = min(max(max_len * 0.12, 0.6), 3.0)
                table.columns[col_idx].width = Inches(width)
            return


def generate_full_report(em_data, sms_data, sl_data, sms_format, campaign_name, template_bytes):
    """
    Generate the complete PPT report. Fully generic — adapts to any campaign data.
    Section order:
      1. Title
      2. Section 1 — Campaign Overview  (Engagement Results table)
      3. Section 2 — Email High-Level Results  (Summary by Touch + Performance Overview)
      4. Section 3 — Subject Line Testing  (one slide per audience, variant rows)
      5. Section 4 — RCM/SMS Engagement  (overview + one slide per Touch)
      6. Section 5 — Creative Engagement  (empty divider for manual heatmap)
    Sections are skipped automatically if no data is available for them.
    """
    DS = _import_functions()

    prs = Presentation(io.BytesIO(template_bytes))
    nav = DS.create_navigation_screen(prs)

    # ===== TITLE SLIDE =====
    prs = DS.duplicate_slide(prs, nav['titleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', campaign_name, verbose=False)

    # ===== SECTION 1: CAMPAIGN OVERVIEW =====
    prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Campaign Overview', verbose=False)

    eng_summary = process_engagement_summary(em_data, sms_data, sms_format)
    if len(eng_summary) > 0:
        prs = DS.duplicate_slide(prs, nav['datatableslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Engagement Results', verbose=False)
        add_data_to_table(prs.slides[-1], eng_summary)

    # ===== SECTION 2: EMAIL HIGH-LEVEL RESULTS =====
    if em_data is not None and len(em_data) > 0:
        prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email High-Level Results', verbose=False)

        em_touch = process_email_summary_by_touch(em_data)
        prs = DS.duplicate_slide(prs, nav['datatableslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Summary by Touch', verbose=False)
        add_data_to_table(prs.slides[-1], em_touch)

        em_detail = process_em_detail(em_data)
        prs = DS.duplicate_slide(prs, nav['datatableslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Email Performance Overview', verbose=False)
        add_data_to_table(prs.slides[-1], em_detail)

    # ===== SECTION 3: SUBJECT LINE TESTING =====
    sl_groups = process_sl_testing(sl_data)
    if sl_groups:
        prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Subject Line Testing', verbose=False)

        for audience_label, sl_table in sl_groups:
            prs = DS.duplicate_slide(prs, nav['datatableslide'], verbose=False)
            title = f"SL Test Results — {truncate_text(audience_label, 40)}"
            DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', title, verbose=False)
            add_data_to_table(prs.slides[-1], sl_table)

    # ===== SECTION 4: RCM / SMS ENGAGEMENT =====
    if sms_data is not None and len(sms_data) > 0:
        label = "RCM/SMS Engagement" if sms_format == "rcm" else "SMS Engagement"
        prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
        DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', label, verbose=False)

        # Overview slide (all touches combined)
        rcm_overview = process_rcm_detail(sms_data, sms_format)
        if len(rcm_overview) > 0:
            prs = DS.duplicate_slide(prs, nav['datatableslide'], verbose=False)
            DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE',
                                 f'{label} Overview', verbose=False)
            add_data_to_table(prs.slides[-1], rcm_overview)

        # One detail slide per Touch
        touches = sms_data['Touch'].unique()
        for touch in touches:
            touch_data = sms_data[sms_data['Touch'] == touch]
            detail = process_rcm_detail(touch_data, sms_format)
            if len(detail) > 0:
                prs = DS.duplicate_slide(prs, nav['datatableslide'], verbose=False)
                slide_title = f'RBM Performance Overview — {touch}' if sms_format == "rcm" \
                    else f'SMS Performance Overview — {touch}'
                DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', slide_title, verbose=False)
                add_data_to_table(prs.slides[-1], detail)

    # ===== SECTION 5: CREATIVE ENGAGEMENT (empty divider) =====
    prs = DS.duplicate_slide(prs, nav['subtitleslide'], verbose=False)
    DS.find_replace_text(prs.slides[-1], 'PLACE_TEXT_TITLE', 'Creative Engagement', verbose=False)

    # ===== CLEANUP template slides =====
    for i in sorted(nav.values(), reverse=True):
        rId = prs.slides._sldIdLst[i]
        prs.slides._sldIdLst.remove(rId)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# =====================================================================
# PULSE CHECK DOC GENERATION
# =====================================================================

def _build_pulse_email_total(em_data):
    """Build Email Summary (Total) table for pulse check."""
    if em_data is None or len(em_data) == 0:
        return None
    total_del = em_data['Deliveries'].sum()
    total_opens = em_data['Unique Opens'].sum()
    total_clicks = em_data['Unique Clicks'].sum()
    return pd.DataFrame([{
        'Metric': 'Total',
        'Deliveries': format_number(total_del),
        'Open Rate': format_percent(_safe_div(total_opens, total_del)),
        'CTR': format_percent(_safe_div(total_clicks, total_del)),
    }])


def _build_pulse_email_by_touch(em_data):
    """Build Email Summary by Touch table for pulse check."""
    return process_email_summary_by_touch(em_data)


def _build_pulse_sms_total(sms_data, sms_format):
    """Build SMS/RCM Summary (Total) table for pulse check."""
    if sms_data is None or len(sms_data) == 0:
        return None
    df = sms_data
    if sms_format == "rcm":
        sms_del = df['SMS Deliveries'].sum() if 'SMS Deliveries' in df.columns else 0
        rbm_del = df['RBM Deliveries'].sum() if 'RBM Deliveries' in df.columns else 0
        sms_click_col = _find_column(df, ['SMS InfoBip Clicks', 'SMS Clicks'])
        rbm_click_col = _find_column(df, ['RBM InfoBip Clicks', 'RBM Clicks'])
        sms_clicks = df[sms_click_col].sum() if sms_click_col else 0
        rbm_clicks = df[rbm_click_col].sum() if rbm_click_col else 0
        return pd.DataFrame([{
            'Metric': 'Total',
            'RBM Deliveries': format_number(rbm_del),
            'SMS Deliveries': format_number(sms_del),
            'RBM CTR': format_percent(_safe_div(rbm_clicks, rbm_del)),
            'SMS CTR': format_percent(_safe_div(sms_clicks, sms_del)),
        }])
    else:
        total_del = df['Deliveries'].sum() if 'Deliveries' in df.columns else 0
        clicks_col = _find_column(df, ['PBI Unique Clicks', 'Branch Unique Clicks', 'Unique Clicks'])
        total_clicks = df[clicks_col].sum() if clicks_col else 0
        return pd.DataFrame([{
            'Metric': 'Total',
            'Deliveries': format_number(total_del),
            'CTR': format_percent(_safe_div(total_clicks, total_del)),
        }])


def _build_pulse_sms_by_touch(sms_data, sms_format):
    """Build SMS/RCM Summary by Touch table for pulse check."""
    if sms_data is None or len(sms_data) == 0:
        return None
    rows = []
    for touch in sms_data['Touch'].unique():
        tdf = sms_data[sms_data['Touch'] == touch]
        if sms_format == "rcm":
            sms_del = tdf['SMS Deliveries'].sum() if 'SMS Deliveries' in tdf.columns else 0
            rbm_del = tdf['RBM Deliveries'].sum() if 'RBM Deliveries' in tdf.columns else 0
            sms_click_col = _find_column(tdf, ['SMS InfoBip Clicks', 'SMS Clicks'])
            rbm_click_col = _find_column(tdf, ['RBM InfoBip Clicks', 'RBM Clicks'])
            sms_clicks = tdf[sms_click_col].sum() if sms_click_col else 0
            rbm_clicks = tdf[rbm_click_col].sum() if rbm_click_col else 0
            rows.append({
                'Touch': touch,
                'RBM Deliveries': format_number(rbm_del),
                'SMS Deliveries': format_number(sms_del),
                'RBM CTR': format_percent(_safe_div(rbm_clicks, rbm_del)),
                'SMS CTR': format_percent(_safe_div(sms_clicks, sms_del)),
            })
        else:
            total_del = tdf['Deliveries'].sum() if 'Deliveries' in tdf.columns else 0
            clicks_col = _find_column(tdf, ['PBI Unique Clicks', 'Branch Unique Clicks', 'Unique Clicks'])
            total_clicks = tdf[clicks_col].sum() if clicks_col else 0
            rows.append({
                'Touch': touch,
                'Deliveries': format_number(total_del),
                'CTR': format_percent(_safe_div(total_clicks, total_del)),
            })
    return pd.DataFrame(rows) if rows else None


def generate_pulse_check(em_data, sms_data, sl_data, sms_format, campaign_name, template_path):
    """
    Generate Pulse Check Word document (.docx) mirroring PPT report content.
    Returns bytes of the generated docx.

    Sections included (when data is available):
      - Email Summary (Total)
      - Email Summary by Touch
      - Email Performance Overview (full detail)
      - Subject Line Testing (one table per audience)
      - RCM/SMS Summary (Total)
      - RCM/SMS Summary by Touch
      - RCM/SMS Performance Overview (full detail)
    """
    sections = []

    # --- Email sections ---
    if em_data is not None and len(em_data) > 0:
        em_total = _build_pulse_email_total(em_data)
        if em_total is not None and len(em_total) > 0:
            sections.append({
                "header": "Email Summary (Total)",
                "data": em_total,
                "widths": [0.25, 0.25, 0.25, 0.25],
            })

        em_touch = _build_pulse_email_by_touch(em_data)
        if em_touch is not None and len(em_touch) > 0:
            sections.append({
                "header": "Email Summary by Touch",
                "data": em_touch,
                "widths": [0.30, 0.23, 0.23, 0.24],
            })

        em_detail = process_em_detail(em_data)
        if em_detail is not None and len(em_detail) > 0:
            n = len(em_detail.columns)
            sections.append({
                "header": "Email Performance Overview",
                "data": em_detail,
                "widths": [1.0 / n] * n,
            })

    # --- Subject Line Testing ---
    sl_groups = process_sl_testing(sl_data)
    if sl_groups:
        for audience_label, sl_table in sl_groups:
            n = len(sl_table.columns)
            sections.append({
                "header": f"SL Testing — {audience_label}",
                "data": sl_table,
                "widths": [1.0 / n] * n,
            })

    # --- SMS / RCM sections ---
    if sms_data is not None and len(sms_data) > 0:
        channel_label = "RCM/SMS" if sms_format == "rcm" else "SMS"

        sms_total = _build_pulse_sms_total(sms_data, sms_format)
        if sms_total is not None and len(sms_total) > 0:
            n = len(sms_total.columns)
            sections.append({
                "header": f"{channel_label} Summary (Total)",
                "data": sms_total,
                "widths": [1.0 / n] * n,
            })

        sms_touch = _build_pulse_sms_by_touch(sms_data, sms_format)
        if sms_touch is not None and len(sms_touch) > 0:
            n = len(sms_touch.columns)
            sections.append({
                "header": f"{channel_label} Summary by Touch",
                "data": sms_touch,
                "widths": [1.0 / n] * n,
            })

        rcm_detail = process_rcm_detail(sms_data, sms_format)
        if rcm_detail is not None and len(rcm_detail) > 0:
            n = len(rcm_detail.columns)
            sections.append({
                "header": f"{channel_label} Performance Overview",
                "data": rcm_detail,
                "widths": [1.0 / n] * n,
            })

    title_line = f"{campaign_name} — Pulse Check"

    return build_pulse_check_docx(
        template_path=template_path,
        title_line=title_line,
        sections=sections,
    )


# =====================================================================
# STREAMLIT UI
# =====================================================================

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, 'template.pptx')
PULSE_CHECK_TEMPLATE_PATH = os.path.join(BASE_DIR, 'pulse_check_template.docx')
SAMPLE_DATA_PATH = os.path.join(BASE_DIR, 'PPT_Sample_Input_Data.xlsx')

# --- Pipeline Overview ---
st.markdown("""
<div class="info-box">
<strong>📋 Auto Load Pipeline:</strong><br>
① Upload cleaned campaign data (Excel with EM / SMS or RCM SMS sheets) →
② Choose output format (PPT Report or Pulse Check Doc) →
③ One-click generate<br>
<em>Supports both standard format (EM/SMS/SLs) and RCM dual-channel format (EM/RCM SMS)</em>
</div>
""", unsafe_allow_html=True)

# --- Data Source Selection ---
st.markdown("---")
st.subheader("📊 Campaign Data")

use_sample_data = st.checkbox("Use built-in sample data", value=False)
if not use_sample_data:
    campaign_file = st.file_uploader(
        "Upload your cleaned campaign data Excel file",
        type=['xlsx', 'xls'],
        help="Excel file with EM, SMS (or RCM SMS), and optionally SLs sheets"
    )
else:
    campaign_file = None

# --- Load and Preview Data ---
st.markdown("---")

loaded = None
load_errors = []

try:
    if use_sample_data and os.path.exists(SAMPLE_DATA_PATH):
        loaded = detect_and_load(SAMPLE_DATA_PATH)
    elif campaign_file is not None:
        loaded = detect_and_load(campaign_file)
except Exception as e:
    load_errors.append(f"Data loading error: {e}")

for err in load_errors:
    st.error(f"❌ {err}")

if loaded:
    em_data = loaded["em_data"]
    sms_data = loaded["sms_data"]
    sl_data = loaded["sl_data"]
    sms_format = loaded["sms_format"]

    # Run validation
    warnings, errors = validate_data(em_data, sms_data, sms_format)
    for w in warnings:
        st.warning(f"⚠️ {w}")
    for e in errors:
        st.error(f"❌ {e}")

    # Summary
    parts = []
    if em_data is not None:
        parts.append(f"{len(em_data)} EM rows")
    if sms_data is not None:
        label = "RCM SMS" if sms_format == "rcm" else "SMS"
        parts.append(f"{len(sms_data)} {label} rows")
    if sl_data is not None:
        parts.append(f"{len(sl_data)} SL Testing rows")

    st.success(f"✅ Data loaded ({sms_format.upper()} format): {', '.join(parts)}")
    st.caption(f"Sheets found: {loaded['sheets_found']}")

    # Data Preview
    st.subheader("👁️ Data Preview")
    tabs = ["📧 Email"]
    if sms_format == "rcm":
        tabs.append("📱 RCM SMS")
    else:
        tabs.append("💬 SMS")
    tabs.append("🔬 Subject Lines")

    tab_list = st.tabs(tabs)

    with tab_list[0]:
        if em_data is not None:
            st.dataframe(em_data, use_container_width=True)
        else:
            st.info("No email data found")

    with tab_list[1]:
        if sms_data is not None:
            st.dataframe(sms_data, use_container_width=True)
        else:
            st.info("No SMS/RCM data found")

    with tab_list[2]:
        if sl_data is not None:
            st.dataframe(sl_data, use_container_width=True)
        else:
            st.info("No subject line data (derived or explicit)")

    # --- Configuration ---
    st.markdown("---")
    st.subheader("⚙️ Configuration")

    # Auto-detect campaign name from data
    default_name = "Campaign Report"
    if em_data is not None and 'Campaign' in em_data.columns:
        campaigns = em_data['Campaign'].dropna().unique()
        if len(campaigns) > 0:
            default_name = str(campaigns[0])
    elif sms_data is not None and 'Campaign' in sms_data.columns:
        campaigns = sms_data['Campaign'].dropna().unique()
        if len(campaigns) > 0:
            default_name = str(campaigns[0])

    campaign_name = st.text_input("Campaign Name (used as report title)", value=default_name)

    output_format = st.radio(
        "Output Format",
        ["📊 Campaign Report (PPT)", "📝 Pulse Check (Word .docx)"],
        horizontal=True,
    )
    is_pulse_check = "Pulse Check" in output_format

    # --- Generate Button ---
    st.markdown("---")

    # Check template availability based on chosen output format
    if is_pulse_check:
        template_ok = os.path.exists(PULSE_CHECK_TEMPLATE_PATH)
        template_missing_msg = "pulse_check_template.docx"
    else:
        template_ok = os.path.exists(TEMPLATE_PATH)
        template_missing_msg = "template.pptx"

    can_generate = (em_data is not None or sms_data is not None) and \
                   template_ok and len(errors) == 0

    button_label = "🚀 Generate Pulse Check (Word)" if is_pulse_check else "🚀 Generate Full Report (PPT)"

    if st.button(button_label, use_container_width=True, disabled=not can_generate):
        with st.spinner("Building your report..."):
            progress = st.progress(0, text="Loading template...")

            if is_pulse_check:
                # ===== PULSE CHECK GENERATION =====
                progress.progress(30, text="Processing data...")
                progress.progress(60, text="Loading Word template...")
                progress.progress(80, text="Generating Pulse Check document...")

                try:
                    docx_bytes = generate_pulse_check(
                        em_data, sms_data, sl_data,
                        sms_format, campaign_name, PULSE_CHECK_TEMPLATE_PATH
                    )
                    progress.progress(100, text="Done!")

                    st.markdown('<div class="success-box">✅ Pulse Check Generated Successfully!</div>',
                                unsafe_allow_html=True)

                    st.info(f"📄 Format: Pulse Check (Word) | "
                            f"Data: {sms_format.upper()} | "
                            f"EM: {len(em_data) if em_data is not None else 0} rows | "
                            f"SMS: {len(sms_data) if sms_data is not None else 0} rows | "
                            f"SL groups: {len(process_sl_testing(sl_data))}")

                    safe_filename = "".join(c for c in campaign_name if c.isalnum() or c in (' ', '-', '_')).strip()
                    safe_filename = safe_filename.replace(' ', '_')
                    st.download_button(
                        label="📥 Download Pulse Check (Word)",
                        data=docx_bytes,
                        file_name=f"{safe_filename}_Pulse_Check.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
                except Exception as e:
                    st.error(f"❌ Generation failed: {e}")
                    import traceback
                    st.code(traceback.format_exc())

            else:
                # ===== PPT GENERATION =====
                progress.progress(30, text="Processing data...")
                progress.progress(60, text="Loading PPT template...")
                with open(TEMPLATE_PATH, 'rb') as f:
                    template_bytes = f.read()

                progress.progress(80, text="Generating PPT slides...")

                try:
                    ppt_output = generate_full_report(
                        em_data, sms_data, sl_data,
                        sms_format, campaign_name, template_bytes
                    )

                    progress.progress(100, text="Done!")

                    st.markdown('<div class="success-box">✅ Report Generated Successfully!</div>',
                                unsafe_allow_html=True)

                    # Slide count summary
                    check_prs = Presentation(io.BytesIO(ppt_output.getvalue()))
                    st.info(f"📄 Total slides: {len(check_prs.slides)} | "
                            f"Format: {sms_format.upper()} | "
                            f"EM: {len(em_data) if em_data is not None else 0} rows | "
                            f"SMS: {len(sms_data) if sms_data is not None else 0} rows | "
                            f"SL groups: {len(process_sl_testing(sl_data))}")

                    safe_filename = "".join(c for c in campaign_name if c.isalnum() or c in (' ', '-', '_')).strip()
                    safe_filename = safe_filename.replace(' ', '_')
                    st.download_button(
                        label="📥 Download Full Campaign Report (PPT)",
                        data=ppt_output,
                        file_name=f"{safe_filename}_report.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                except Exception as e:
                    st.error(f"❌ Generation failed: {e}")
                    import traceback
                    st.code(traceback.format_exc())

    elif not can_generate:
        if len(errors) > 0:
            st.warning("⚠️ Please fix data errors above before generating.")
        elif not template_ok:
            st.warning(f"⚠️ Template file missing: {template_missing_msg}. Please ensure it exists in the project folder.")
        else:
            st.warning("⚠️ Please load campaign data to generate report.")

else:
    if not load_errors:
        st.info("👆 Upload your cleaned campaign data Excel file above, or check the sample data box.")

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #888; font-size: 0.9rem;">
    🚀 Auto Load Demo | Postpaid Marketing Analytics — Campaign Report Generator
</div>
""", unsafe_allow_html=True)
